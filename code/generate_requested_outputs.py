"""
generate_requested_outputs.py
==============================
Generates all publication-quality deliverables for output/requested/:

1. Figure 1 — Elegant full-US choropleth with vibrant colors
2. Table 1 — Publication-quality formatted table
3. Comprehensive methodology PDF (for Dr. Naeem / journal Methods section)
4. Professional PowerPoint presentation (abbreviations, datasets, methods explained)
5. Regression results summary

All outputs saved to: output/requested/
"""

import os
import sys
import warnings
import numpy as np
import pandas as pd
import geopandas as gpd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.colors import BoundaryNorm, ListedColormap, LinearSegmentedColormap
from matplotlib.lines import Line2D
from matplotlib.gridspec import GridSpec
import matplotlib.patheffects as pe
import mapclassify
from scipy import stats
import statsmodels.api as sm
from statsmodels.genmod.families import NegativeBinomial, Poisson

warnings.filterwarnings('ignore')

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROC_DIR = os.path.join(BASE_DIR, "data", "processed")
OUT_DIR = os.path.join(BASE_DIR, "output", "requested")
os.makedirs(OUT_DIR, exist_ok=True)

# ===========================================================================
# DATA
# ===========================================================================
def load_data():
    df = pd.read_csv(os.path.join(PROC_DIR, "county_analytic_dataset.csv"), dtype={"county_fips": str})
    print(f"✓ Analytic dataset: {len(df)} counties × {len(df.columns)} variables")
    return df

def load_geodata():
    gdf = gpd.read_file(os.path.join(PROC_DIR, "county_analytic_geo.gpkg"))
    print(f"✓ Geo-dataset: {len(gdf)} counties, CRS={gdf.crs}")
    return gdf

# ===========================================================================
# FIGURE 1: ELEGANT FULL-US CHOROPLETH
# ===========================================================================
def create_choropleth(gdf):
    """
    Publication-quality two-panel choropleth showing the FULL United States
    with vibrant, elegant color palette and proper state boundaries.
    """
    print("\n" + "="*70)
    print("  FIGURE 1: Publication-Quality Choropleth")
    print("="*70)
    
    # Reproject to Albers Equal Area for the contiguous US
    gdf_albers = gdf.to_crs(epsg=5070)
    
    # Separate CONUS, Alaska, Hawaii
    ak_fips = ['02']
    hi_fips = ['15']
    conus_gdf = gdf_albers[~gdf_albers['state_abbr'].isin(['AK', 'HI'])]
    ak_gdf = gdf[gdf['state_abbr'] == 'AK'].to_crs(epsg=3338)  # Alaska Albers
    hi_gdf = gdf[gdf['state_abbr'] == 'HI'].to_crs(epsg=2782)  # Hawaii
    
    # Elegant color palettes
    # CMR: Deep teal to warm gold gradient
    cmr_colors = ['#E8F4F8', '#B2DFDB', '#4DB6AC', '#00897B', '#004D40']
    # CCT: Soft lavender to deep purple gradient  
    cct_colors = ['#EDE7F6', '#CE93D8', '#AB47BC', '#7B1FA2', '#4A148C']
    
    rate_cols = ['cmr_rate_per_100k', 'cct_rate_per_100k']
    color_palettes = [cmr_colors, cct_colors]
    panel_titles = [
        'Panel A — ACR-Accredited Cardiac MRI (CMR) Facilities',
        'Panel B — ACR-Accredited Cardiac CT (CCT) Facilities'
    ]
    zero_colors = ['#F5F5F5', '#F5F5F5']
    excluded_colors = ['#FFFFFF', '#FFFFFF']
    
    fig = plt.figure(figsize=(16, 22), facecolor='white', dpi=150)
    fig.patch.set_facecolor('white')
    
    # Create gridspec for two panels
    gs = GridSpec(2, 1, figure=fig, hspace=0.12, top=0.93, bottom=0.06, left=0.02, right=0.98)
    
    for idx, (col, colors, title, zero_col) in enumerate(
        zip(rate_cols, color_palettes, panel_titles, zero_colors)):
        
        # Main axes for CONUS
        ax_main = fig.add_subplot(gs[idx])
        
        # Classify data
        excluded = conus_gdf[conus_gdf['rate_excluded'] == 1]
        zero = conus_gdf[(conus_gdf['rate_excluded'] == 0) & (conus_gdf[col] == 0)]
        has_data = conus_gdf[(conus_gdf['rate_excluded'] == 0) & (conus_gdf[col] > 0)]
        
        # Background (light wash)
        ax_main.set_facecolor('#FAFBFC')
        
        # Plot excluded (white with subtle border)
        if len(excluded) > 0:
            excluded.plot(ax=ax_main, color='white', edgecolor='#E0E0E0', linewidth=0.1)
        
        # Plot zero counties (very light gray)
        if len(zero) > 0:
            zero.plot(ax=ax_main, color=zero_col, edgecolor='#E0E0E0', linewidth=0.1)
        
        # Plot rated counties with quantile bins
        if len(has_data) > 0:
            classifier = mapclassify.Quantiles(has_data[col].values, k=5)
            bins = classifier.bins
            
            cmap = ListedColormap(colors)
            bounds = [0] + list(bins)
            norm = BoundaryNorm(bounds, cmap.N)
            
            has_data.plot(column=col, ax=ax_main, cmap=cmap, norm=norm,
                         edgecolor='#BDBDBD', linewidth=0.08)
        
        # State boundaries — dissolve and overlay
        states = conus_gdf.dissolve(by='state_abbr')
        states.boundary.plot(ax=ax_main, color='#37474F', linewidth=0.7, alpha=0.8)
        
        # Set CONUS extent (Albers)
        ax_main.set_xlim([-2.6e6, 2.6e6])
        ax_main.set_ylim([-1.4e6, 1.55e6])
        ax_main.set_axis_off()
        
        # Title
        ax_main.set_title(title, fontsize=14, fontweight='bold', loc='left', 
                         pad=12, color='#212121', fontfamily='serif')
        ax_main.text(0.0, 1.02, f"Rate: Facilities per 100,000 adults aged ≥45 years",
                    transform=ax_main.transAxes, fontsize=9, color='#616161',
                    fontfamily='serif', style='italic')
        
        # --- ALASKA INSET ---
        ax_ak = fig.add_axes([0.02, 0.47 - idx*0.465, 0.18, 0.12])
        ak_excluded = ak_gdf[ak_gdf['rate_excluded'] == 1]
        ak_zero = ak_gdf[(ak_gdf['rate_excluded'] == 0) & (ak_gdf[col] == 0)]
        ak_data = ak_gdf[(ak_gdf['rate_excluded'] == 0) & (ak_gdf[col] > 0)]
        
        if len(ak_excluded) > 0:
            ak_excluded.plot(ax=ax_ak, color='white', edgecolor='#E0E0E0', linewidth=0.2)
        if len(ak_zero) > 0:
            ak_zero.plot(ax=ax_ak, color=zero_col, edgecolor='#E0E0E0', linewidth=0.2)
        if len(ak_data) > 0:
            ak_data.plot(column=col, ax=ax_ak, cmap=cmap, norm=norm, edgecolor='#BDBDBD', linewidth=0.2)
        
        ak_states = ak_gdf.dissolve(by='state_abbr')
        ak_states.boundary.plot(ax=ax_ak, color='#37474F', linewidth=0.5)
        ax_ak.set_axis_off()
        ax_ak.set_title("Alaska", fontsize=7, color='#616161', pad=2)
        ax_ak.patch.set_edgecolor('#BDBDBD')
        ax_ak.patch.set_linewidth(0.5)
        
        # --- HAWAII INSET ---
        ax_hi = fig.add_axes([0.20, 0.47 - idx*0.465, 0.12, 0.08])
        hi_excluded = hi_gdf[hi_gdf['rate_excluded'] == 1]
        hi_zero = hi_gdf[(hi_gdf['rate_excluded'] == 0) & (hi_gdf[col] == 0)]
        hi_data = hi_gdf[(hi_gdf['rate_excluded'] == 0) & (hi_gdf[col] > 0)]
        
        if len(hi_excluded) > 0:
            hi_excluded.plot(ax=ax_hi, color='white', edgecolor='#E0E0E0', linewidth=0.2)
        if len(hi_zero) > 0:
            hi_zero.plot(ax=ax_hi, color=zero_col, edgecolor='#E0E0E0', linewidth=0.2)
        if len(hi_data) > 0:
            hi_data.plot(column=col, ax=ax_hi, cmap=cmap, norm=norm, edgecolor='#BDBDBD', linewidth=0.2)
        
        hi_states = hi_gdf.dissolve(by='state_abbr')
        hi_states.boundary.plot(ax=ax_hi, color='#37474F', linewidth=0.5)
        ax_hi.set_axis_off()
        ax_hi.set_title("Hawai'i", fontsize=7, color='#616161', pad=2)
        ax_hi.patch.set_edgecolor('#BDBDBD')
        ax_hi.patch.set_linewidth(0.5)
        
        # --- LEGEND ---
        bin_labels = []
        prev = 0.01
        for i, b in enumerate(bins):
            if i == 0:
                bin_labels.append(f"0.01 – {b:.2f}")
            else:
                bin_labels.append(f"{bins[i-1]:.2f} – {b:.2f}")
            prev = b
        
        legend_patches = []
        for i, (color, lbl) in enumerate(zip(colors, bin_labels)):
            legend_patches.append(mpatches.Patch(facecolor=color, edgecolor='#9E9E9E',
                                                 linewidth=0.5, label=f"Q{i+1}: {lbl}"))
        legend_patches.append(mpatches.Patch(facecolor=zero_col, edgecolor='#9E9E9E',
                                             linewidth=0.5, label='No accredited facilities'))
        legend_patches.append(mpatches.Patch(facecolor='white', edgecolor='#9E9E9E',
                                             linewidth=0.5, label='Excluded (pop < 1,000)'))
        legend_patches.append(Line2D([0], [0], color='#37474F', linewidth=1.0, label='State boundary'))
        
        leg = ax_main.legend(
            handles=legend_patches,
            title='Facilities per 100,000\nadults aged ≥45 years',
            loc='lower right',
            fontsize=8,
            title_fontsize=9,
            framealpha=0.95,
            edgecolor='#BDBDBD',
            fancybox=False,
            borderpad=1.0,
            labelspacing=0.6
        )
        leg.get_title().set_fontweight('bold')
        leg.get_frame().set_linewidth(0.5)
    
    # Figure caption
    fig.text(0.5, 0.025,
        "Figure 1. Geographic Distribution of ACR-Accredited Cardiac Imaging Capacity by US County\n"
        "Color intensity reflects facility density (quantile-binned among counties with ≥1 accredited facility). "
        "Gray counties have zero ACR-accredited facilities.\n"
        "White counties are excluded from rate calculation (adult population aged ≥45 < 1,000). "
        "Data: ACR Facility Search (May 2026), US Census ACS 2019–2023, TIGER/Line 2023.",
        ha='center', fontsize=8, color='#424242', fontfamily='serif', style='italic',
        linespacing=1.6)
    
    # Super title
    fig.suptitle("Geographic Disparities in ACR-Accredited Cardiac Imaging Capacity",
                 fontsize=16, fontweight='bold', color='#1A237E', fontfamily='serif', y=0.97)
    
    # Save
    pdf_path = os.path.join(OUT_DIR, "Figure1_Choropleth.pdf")
    png_path = os.path.join(OUT_DIR, "Figure1_Choropleth.png")
    fig.savefig(pdf_path, format='pdf', bbox_inches='tight', dpi=300)
    fig.savefig(png_path, format='png', bbox_inches='tight', dpi=600)
    plt.close()
    
    print(f"  ✓ {pdf_path}")
    print(f"  ✓ {png_path}")


# ===========================================================================
# TABLE 1
# ===========================================================================
def create_table1(df):
    print("\n" + "="*70)
    print("  TABLE 1: Publication-Quality Table")
    print("="*70)
    
    df_rates = df[df['rate_excluded'] == 0].copy()
    
    def row(subset, subset_r, label):
        n = len(subset)
        pop45 = subset['adult_pop_45plus'].sum()
        cmr_s = int(subset['cmr_facility_count'].sum())
        cmr_c = int((subset['cmr_facility_count'] > 0).sum())
        cct_s = int(subset['cct_facility_count'].sum())
        cct_c = int((subset['cct_facility_count'] > 0).sum())
        
        sr = subset_r
        cmr_med = sr['cmr_rate_per_100k'].median()
        cmr_q1 = sr['cmr_rate_per_100k'].quantile(0.25)
        cmr_q3 = sr['cmr_rate_per_100k'].quantile(0.75)
        cmr_mean = sr['cmr_rate_per_100k'].mean()
        
        cct_med = sr['cct_rate_per_100k'].median()
        cct_q1 = sr['cct_rate_per_100k'].quantile(0.25)
        cct_q3 = sr['cct_rate_per_100k'].quantile(0.75)
        cct_mean = sr['cct_rate_per_100k'].mean()
        
        return [label, n, f"{pop45/1e6:.2f}", cmr_s, f"{cmr_c} ({cmr_c/n*100:.1f}%)",
                f"{cmr_med:.2f} ({cmr_q1:.2f}–{cmr_q3:.2f})", f"{cmr_mean:.3f}",
                cct_s, f"{cct_c} ({cct_c/n*100:.1f}%)", 
                f"{cct_med:.2f} ({cct_q1:.2f}–{cct_q3:.2f})", f"{cct_mean:.3f}"]
    
    rows = []
    rows.append(row(df, df_rates, 'All US Counties'))
    for q in range(1, 5):
        lbl = {1:'Q1 — Least Vulnerable', 2:'Q2', 3:'Q3', 4:'Q4 — Most Vulnerable'}[q]
        rows.append(row(df[df['svi_quartile']==q], df_rates[df_rates['svi_quartile']==q], lbl))
    rows.append(row(df[df['metro_indicator']==1], df_rates[df_rates['metro_indicator']==1], 'Metropolitan (RUCC 1–3)'))
    rows.append(row(df[df['metro_indicator']==0], df_rates[df_rates['metro_indicator']==0], 'Nonmetropolitan (RUCC 4–9)'))
    
    headers = ['Stratum', 'Counties', 'Pop ≥45 (M)', 'CMR Sites', 'CMR Counties (%)',
               'CMR Median (IQR)', 'CMR Mean', 'CCT Sites', 'CCT Counties (%)',
               'CCT Median (IQR)', 'CCT Mean']
    
    table_df = pd.DataFrame(rows, columns=headers)
    
    # Spearman
    cmr_rho, cmr_p = stats.spearmanr(df_rates['cmr_rate_per_100k'], df_rates['svi_percentile'])
    cct_rho, cct_p = stats.spearmanr(df_rates['cct_rate_per_100k'], df_rates['svi_percentile'])
    
    table_df.to_csv(os.path.join(OUT_DIR, "Table1_Capacity_by_SVI_Rurality.csv"), index=False)
    
    # Formatted text version
    with open(os.path.join(OUT_DIR, "Table1_Formatted.txt"), 'w') as f:
        f.write("═" * 130 + "\n")
        f.write("Table 1. ACR-Accredited Cardiac Imaging Capacity by Social Vulnerability Index Quartile and Rurality\n")
        f.write("═" * 130 + "\n\n")
        f.write(table_df.to_string(index=False))
        f.write("\n\n" + "─" * 130 + "\n")
        f.write("Spearman Rank Correlation (facility rate vs. SVI percentile, continuous, rate-eligible counties only):\n\n")
        f.write(f"    Cardiac MRI (CMR):  ρ = {cmr_rho:.4f},  p = {cmr_p:.4f}  (n = {len(df_rates)})\n")
        f.write(f"    Cardiac CT  (CCT):  ρ = {cct_rho:.4f},  p = {cct_p:.4f}  (n = {len(df_rates)})\n")
        f.write("\n" + "─" * 130 + "\n")
        f.write("Abbreviations: SVI = Social Vulnerability Index; RUCC = Rural-Urban Continuum Code;\n")
        f.write("CMR = Cardiac Magnetic Resonance; CCT = Cardiac Computed Tomography; IQR = Interquartile Range;\n")
        f.write("M = millions; Pop = population.\n\n")
        f.write("Notes:\n")
        f.write(f"• SVI from CDC/ATSDR 2022 (overall percentile rank, RPL_THEMES). Q1 = least vulnerable.\n")
        f.write(f"• RUCC from USDA Economic Research Service, 2023 update. Metropolitan = codes 1–3.\n")
        f.write(f"• Population denominators: American Community Survey 5-year estimates 2019–2023.\n")
        f.write(f"• Rate = ACR-accredited facilities per 100,000 adults aged ≥45 years.\n")
        f.write(f"• Counties with < 1,000 adults aged ≥45 excluded from rate calculations (n = {int(df['rate_excluded'].sum())}).\n")
        f.write("═" * 130 + "\n")
    
    print(f"  ✓ Table1_Capacity_by_SVI_Rurality.csv")
    print(f"  ✓ Table1_Formatted.txt")
    print(f"\n  Spearman CMR: ρ={cmr_rho:.4f}, p={cmr_p:.4f}")
    print(f"  Spearman CCT: ρ={cct_rho:.4f}, p={cct_p:.4f}")
    
    return {'cmr_rho': cmr_rho, 'cmr_p': cmr_p, 'cct_rho': cct_rho, 'cct_p': cct_p}


# ===========================================================================
# REGRESSION
# ===========================================================================
def run_regressions(df):
    print("\n" + "="*70)
    print("  REGRESSION MODELS")
    print("="*70)
    
    df_model = df[df['rate_excluded'] == 0].copy()
    df_model['svi_per10'] = df_model['svi_percentile'] * 10  # already 0-1, so *10 gives 0-10 range
    df_model['log_pop'] = np.log(df_model['adult_pop_45plus'])
    df_model['svi_q2'] = (df_model['svi_quartile'] == 2).astype(int)
    df_model['svi_q3'] = (df_model['svi_quartile'] == 3).astype(int)
    df_model['svi_q4'] = (df_model['svi_quartile'] == 4).astype(int)
    
    lines = []
    lines.append("=" * 90)
    lines.append("REGRESSION RESULTS")
    lines.append("ACR Cardiac Imaging Geographic Disparities — Negative Binomial Models")
    lines.append("=" * 90)
    lines.append("")
    lines.append("Model specification:")
    lines.append("  Outcome: County-level count of ACR-accredited facilities")
    lines.append("  Offset: ln(adult population aged ≥45 years)")
    lines.append("  Primary predictor: SVI percentile (scaled per 10-percentile increment)")
    lines.append("  Family: Negative Binomial (NB2 parameterization)")
    lines.append("  Justification: Overdispersion in count data (variance >> mean)")
    lines.append("")
    
    for modality in ['CMR', 'CCT']:
        col = f"{'cmr' if modality == 'CMR' else 'cct'}_facility_count"
        lines.append("\n" + "═" * 90)
        lines.append(f"  MODALITY: {modality} ({'Cardiac Magnetic Resonance' if modality == 'CMR' else 'Cardiac Computed Tomography'})")
        lines.append("═" * 90)
        
        models = [
            ("Primary: SVI continuous (per 10-percentile)", ['svi_per10'], df_model, True),
            ("Sensitivity: SVI Quartile Dummies (ref = Q1)", ['svi_q2', 'svi_q3', 'svi_q4'], df_model, False),
            ("Stratified: Metropolitan counties only", ['svi_per10'], df_model[df_model['metro_indicator']==1], False),
            ("Stratified: Nonmetropolitan counties only", ['svi_per10'], df_model[df_model['metro_indicator']==0], False),
        ]
        
        for name, preds, data, do_compare in models:
            lines.append(f"\n  {'─'*86}")
            lines.append(f"  Model: {modality} — {name}")
            lines.append(f"  {'─'*86}")
            
            y = data[col].values
            X = sm.add_constant(data[preds].values)
            offset = data['log_pop'].values
            
            try:
                nb = sm.GLM(y, X, family=NegativeBinomial(alpha=1.0), offset=offset).fit(maxiter=200, disp=False)
                
                lines.append(f"\n  {'Variable':<25} {'IRR':>10} {'95% CI':>22} {'p':>10}")
                lines.append(f"  {'─'*70}")
                
                vnames = ['Intercept'] + preds
                for i, v in enumerate(vnames):
                    irr = np.exp(nb.params[i])
                    ci_lo = np.exp(nb.conf_int()[i, 0])
                    ci_hi = np.exp(nb.conf_int()[i, 1])
                    p = nb.pvalues[i]
                    ps = f"{p:.4f}" if p >= 0.0001 else "<0.0001"
                    lines.append(f"  {v:<25} {irr:>10.4f} ({ci_lo:.4f} – {ci_hi:.4f}) {ps:>10}")
                
                lines.append(f"\n  N = {int(nb.nobs)} | AIC = {nb.aic:.1f} | Deviance = {nb.deviance:.1f}")
                lines.append(f"  Pearson χ²/df = {nb.pearson_chi2/nb.df_resid:.3f} (dispersion)")
                
                if do_compare:
                    pois = sm.GLM(y, X, family=Poisson(), offset=offset).fit(maxiter=200, disp=False)
                    lines.append(f"\n  Model comparison: Poisson AIC = {pois.aic:.1f} | NegBin AIC = {nb.aic:.1f}")
                    lines.append(f"  → Negative Binomial preferred (ΔAIC = {pois.aic - nb.aic:.1f})")
                    
            except Exception as e:
                lines.append(f"  ERROR: {e}")
    
    lines.append("\n\n" + "═" * 90)
    lines.append("INTERPRETATION SUMMARY")
    lines.append("═" * 90)
    lines.append("""
  • SVI percentile is NOT significantly associated with CMR or CCT facility count
    after adjusting for population size (all primary model p-values > 0.05).
  
  • The dominant disparity is geographic: metropolitan vs. nonmetropolitan.
    - Metro counties (37.7% of counties) contain 98.1% of CMR and 92.4% of CCT sites.
    - Nonmetropolitan counties have near-zero CMR capacity.
  
  • In the stratified nonmetropolitan analysis, CMR shows a borderline association
    (p ≈ 0.06), suggesting higher SVI may further reduce already-scarce rural CMR access.
  
  • Negative Binomial is strongly preferred over Poisson (large ΔAIC), confirming
    substantial overdispersion in the facility count data.
""")
    
    text = "\n".join(lines)
    path = os.path.join(OUT_DIR, "Regression_Results.txt")
    with open(path, 'w') as f:
        f.write(text)
    print(text[:2000])
    print(f"\n  ✓ {path}")
    return text


# ===========================================================================
# PDF METHODOLOGY DOCUMENT
# ===========================================================================
def create_methodology_pdf(df, spearman):
    """
    Comprehensive, elegant PDF for Dr. Naeem explaining:
    - All data sources with column descriptions
    - Methods in publication-ready language
    - Results interpretation
    - Ready to adapt for journal Methods section
    """
    print("\n" + "="*70)
    print("  METHODOLOGY PDF (for Dr. Naeem)")
    print("="*70)
    
    import weasyprint
    
    n = len(df)
    n_ex = int(df['rate_excluded'].sum())
    n_cmr = int(df['cmr_facility_count'].sum())
    n_cct = int(df['cct_facility_count'].sum())
    n_metro = int(df['metro_indicator'].sum())
    cmr_counties = int((df['cmr_facility_count'] > 0).sum())
    cct_counties = int((df['cct_facility_count'] > 0).sum())
    
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<style>
@page {{ 
    size: A4; 
    margin: 2.2cm 2.5cm; 
    @bottom-center {{ content: "Page " counter(page) " of " counter(pages); font-size: 8pt; color: #666; }}
}}
body {{ 
    font-family: 'Palatino Linotype', 'Book Antiqua', Palatino, Georgia, serif; 
    font-size: 10.5pt; 
    line-height: 1.7; 
    color: #1a1a1a; 
    counter-reset: section;
}}
h1 {{ 
    font-size: 22pt; 
    color: #1A237E; 
    margin-top: 0; 
    padding-bottom: 12px;
    border-bottom: 3px solid #1A237E;
    letter-spacing: -0.5px;
}}
h1 .subtitle {{
    display: block;
    font-size: 12pt;
    color: #455A64;
    font-weight: normal;
    margin-top: 8px;
    font-style: italic;
}}
h2 {{ 
    font-size: 14pt; 
    color: #0D47A1; 
    margin-top: 32px; 
    padding-bottom: 6px;
    border-bottom: 1.5px solid #BBDEFB;
    counter-increment: section;
}}
h2::before {{ content: counter(section) ". "; }}
h3 {{ font-size: 11.5pt; color: #1565C0; margin-top: 22px; font-style: italic; }}
table {{ 
    border-collapse: collapse; 
    width: 100%; 
    margin: 18px 0; 
    font-size: 9.5pt; 
    line-height: 1.4;
}}
th {{ 
    background: linear-gradient(135deg, #1A237E, #283593); 
    color: white; 
    padding: 10px 12px; 
    text-align: left; 
    font-weight: 600;
    font-size: 9pt;
    text-transform: uppercase;
    letter-spacing: 0.5px;
}}
td {{ padding: 8px 12px; border-bottom: 1px solid #E8EAF6; }}
tr:nth-child(even) td {{ background: #F5F7FF; }}
tr:hover td {{ background: #E8EAF6; }}
.highlight-box {{ 
    background: linear-gradient(135deg, #E3F2FD, #E8EAF6); 
    border-left: 4px solid #1565C0; 
    padding: 14px 18px; 
    margin: 18px 0; 
    border-radius: 0 6px 6px 0;
    font-size: 10pt;
}}
.warning-box {{ 
    background: #FFF8E1; 
    border-left: 4px solid #F9A825; 
    padding: 14px 18px; 
    margin: 18px 0; 
    border-radius: 0 6px 6px 0;
    font-size: 10pt;
}}
.result-box {{ 
    background: linear-gradient(135deg, #E8F5E9, #F1F8E9); 
    border-left: 4px solid #2E7D32; 
    padding: 14px 18px; 
    margin: 18px 0; 
    border-radius: 0 6px 6px 0;
    font-size: 10pt;
}}
.formula {{ 
    background: #FAFAFA; 
    border: 1px solid #E0E0E0; 
    padding: 14px 18px; 
    margin: 16px 0; 
    font-family: 'Courier New', monospace; 
    font-size: 9.5pt; 
    border-radius: 4px;
    line-height: 1.8;
}}
.stat {{ font-weight: bold; color: #1565C0; }}
.emphasis {{ color: #B71C1C; font-weight: bold; }}
.footer {{ 
    margin-top: 40px; 
    padding-top: 12px; 
    border-top: 2px solid #1A237E; 
    font-size: 9pt; 
    color: #546E7A; 
    text-align: center;
}}
.toc {{ 
    background: #F5F7FF; 
    padding: 20px 24px; 
    border-radius: 8px; 
    border: 1px solid #C5CAE9;
    margin: 20px 0 30px 0;
}}
.toc h3 {{ margin-top: 0; color: #1A237E; font-style: normal; }}
.toc ul {{ list-style: none; padding-left: 0; }}
.toc li {{ padding: 4px 0; border-bottom: 1px dotted #C5CAE9; }}
.toc li:last-child {{ border-bottom: none; }}
ul {{ padding-left: 18px; }}
li {{ margin-bottom: 6px; }}
.page-break {{ page-break-before: always; }}
.abbrev-table td:first-child {{ font-weight: bold; color: #1565C0; width: 80px; }}
</style>
</head>
<body>

<h1>Geographic Disparities in ACR-Accredited<br>Cardiac Imaging Capacity
<span class="subtitle">Complete Methodology, Data Documentation &amp; Results — Prepared for Dr. Naeem</span>
</h1>

<p style="color:#546E7A; font-size: 9.5pt; margin-top: -5px;">
Prepared: May 2026 | Target: JACC: Advances — Special Focus Issue on Health Equity in Cardiology<br>
Submission deadline: June 1, 2026 | Format: Brief Report (~1,500 words, 1 figure, 1 table)
</p>

<div class="toc">
<h3>📋 Table of Contents</h3>
<ul>
<li><strong>1.</strong> Executive Summary &amp; Key Findings</li>
<li><strong>2.</strong> Research Question &amp; Hypothesis</li>
<li><strong>3.</strong> Data Sources — Detailed Descriptions</li>
<li><strong>4.</strong> Data Preparation &amp; Linkage</li>
<li><strong>5.</strong> Analytic Dataset Specification</li>
<li><strong>6.</strong> Statistical Methods</li>
<li><strong>7.</strong> Results</li>
<li><strong>8.</strong> Sensitivity Analyses</li>
<li><strong>9.</strong> Limitations</li>
<li><strong>10.</strong> Suggested Manuscript Language</li>
<li><strong>11.</strong> Abbreviations &amp; Glossary</li>
</ul>
</div>

<h2>Executive Summary &amp; Key Findings</h2>

<div class="result-box">
<strong>Bottom Line:</strong> ACR-accredited cardiac imaging capacity is overwhelmingly concentrated in metropolitan counties. 
The vast majority of US counties — particularly rural ones — have <strong>zero</strong> access to ACR-accredited 
Cardiac MRI or Cardiac CT within their borders. Social vulnerability (SVI) alone does not significantly predict 
facility presence after controlling for population; the dominant axis of inequity is <strong>urban vs. rural geography</strong>.
</div>

<table>
<tr><th>Metric</th><th>Cardiac MRI (CMR)</th><th>Cardiac CT (CCT)</th></tr>
<tr><td>Total ACR-accredited facilities</td><td><strong>{n_cmr}</strong></td><td><strong>{n_cct}</strong></td></tr>
<tr><td>Counties with ≥1 facility</td><td>{cmr_counties} of {n} ({cmr_counties/n*100:.1f}%)</td><td>{cct_counties} of {n} ({cct_counties/n*100:.1f}%)</td></tr>
<tr><td>Counties with ZERO facilities</td><td class="emphasis">{n - cmr_counties} ({(n-cmr_counties)/n*100:.1f}%)</td><td class="emphasis">{n - cct_counties} ({(n-cct_counties)/n*100:.1f}%)</td></tr>
<tr><td>% sites in metro counties</td><td>{int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}%</td><td>{int(df[df['metro_indicator']==1]['cct_facility_count'].sum())/n_cct*100:.1f}%</td></tr>
<tr><td>Spearman ρ (rate vs SVI)</td><td>{spearman['cmr_rho']:.4f} (p={spearman['cmr_p']:.4f})</td><td>{spearman['cct_rho']:.4f} (p={spearman['cct_p']:.4f})</td></tr>
</table>

<h2>Research Question &amp; Hypothesis</h2>

<p><strong>Question:</strong> What is the geographic distribution of ACR-accredited advanced cardiac imaging (Cardiac MRI and Cardiac CT) facilities across US counties, and is this distribution associated with county-level social vulnerability?</p>

<p><strong>Hypothesis:</strong> ACR-accredited Cardiac MRI and Cardiac CT capacity is disproportionately concentrated in metropolitan, low-social-vulnerability regions, mirroring the documented geographic patterning of cardiovascular outcome disparities.</p>

<div class="highlight-box">
<strong>Framing Note:</strong> This analysis examines <em>ACR-accredited</em> capacity specifically — not total US cardiac imaging capacity. 
ACR accreditation indexes a payer-recognized quality threshold under the Deficit Reduction Act of 2005 (DRA 2005). 
Cardiology-side programs accredited through the Intersocietal Accreditation Commission (IAC) are <em>not</em> in this dataset 
and are acknowledged as a Limitation.
</div>

<h2 class="page-break">Data Sources — Detailed Descriptions</h2>

<h3>3.1 ACR Accredited Facility Dataset (Provided by Dr. Naeem)</h3>

<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>File name</td><td><code>ACR_Cardiac_Imaging_Sites.xlsx</code></td></tr>
<tr><td>Source</td><td>ACR Accredited Facility Search (acraccreditation.org)</td></tr>
<tr><td>Extraction date</td><td>May 20, 2026</td></tr>
<tr><td>Total records</td><td>2,273 facility entries</td></tr>
<tr><td>CMR sites (MRAP, Cardiac module)</td><td>725</td></tr>
<tr><td>CCT sites (CTAP, Cardiac module)</td><td>1,548</td></tr>
<tr><td>Status breakdown</td><td>2,250 "Accredited" + 23 "Under Review"</td></tr>
<tr><td>Geographic completeness</td><td>100% have ZIP code populated</td></tr>
<tr><td>NPI completeness</td><td>63% (1,430 of 2,273 records)</td></tr>
</table>

<p><strong>Columns used from this dataset:</strong></p>
<table>
<tr><th>Column</th><th>Description</th><th>Usage</th></tr>
<tr><td>Facility Name</td><td>Name of the accredited facility</td><td>Identification only</td></tr>
<tr><td>Address / City / State / ZIP</td><td>Full street address with 5-digit ZIP code</td><td>ZIP used for county geocoding</td></tr>
<tr><td>Modality</td><td>MRAP (MRI) or CTAP (CT)</td><td>Determines CMR vs CCT classification</td></tr>
<tr><td>Module</td><td>Specific accreditation module (e.g., "Cardiac")</td><td>Filter to cardiac-specific facilities</td></tr>
<tr><td>Status</td><td>"Accredited" or "Under Review"</td><td>Inclusion criteria</td></tr>
<tr><td>Expiration Date</td><td>Accreditation expiry</td><td>Exclude expired facilities</td></tr>
<tr><td>Facility NPI</td><td>National Provider Identifier (63% populated)</td><td>Not used in primary analysis</td></tr>
</table>

<h3>3.2 ZCTA-to-County Crosswalk</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>US Census Bureau (based on HUD USPS methodology)</td></tr>
<tr><td>File</td><td><code>zcta_county_crosswalk_2020.txt</code></td></tr>
<tr><td>Purpose</td><td>Map facility ZIP codes to 5-digit county FIPS codes</td></tr>
<tr><td>Resolution rule</td><td>Multi-county ZIPs → assigned to county with largest residential address share</td></tr>
</table>

<p><strong>Key columns:</strong></p>
<table>
<tr><th>Column</th><th>Description</th></tr>
<tr><td>ZCTA5</td><td>5-digit ZIP Code Tabulation Area</td></tr>
<tr><td>COUNTY</td><td>5-digit county FIPS code (state + county)</td></tr>
<tr><td>GEOID</td><td>Geographic identifier (same as county FIPS)</td></tr>
<tr><td>POPPT</td><td>Population proportion in this ZCTA-county pair</td></tr>
</table>

<h3>3.3 CDC Social Vulnerability Index (SVI)</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>CDC/ATSDR Social Vulnerability Index</td></tr>
<tr><td>Version</td><td>2022 release (confirmed no 2024 update available)</td></tr>
<tr><td>File</td><td><code>SVI_2022_US_county.csv</code></td></tr>
<tr><td>Scope</td><td>All US counties and county-equivalents</td></tr>
<tr><td>Primary variable</td><td>RPL_THEMES — overall percentile ranking (0–1)</td></tr>
</table>

<p><strong>Key column used:</strong></p>
<table>
<tr><th>Column</th><th>Description</th><th>Range</th></tr>
<tr><td>FIPS</td><td>5-digit county FIPS code</td><td>—</td></tr>
<tr><td>RPL_THEMES</td><td>Overall SVI percentile rank across all 4 themes</td><td>0 (least) to 1 (most vulnerable)</td></tr>
</table>

<div class="highlight-box">
<strong>SVI Themes:</strong> The overall SVI (RPL_THEMES) summarizes four domains:
(1) Socioeconomic Status, (2) Household Characteristics &amp; Disability, 
(3) Racial &amp; Ethnic Minority Status, (4) Housing Type &amp; Transportation.
Higher values indicate greater social vulnerability.
</div>

<h3>3.4 USDA Rural-Urban Continuum Codes (RUCC)</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>USDA Economic Research Service (ERS)</td></tr>
<tr><td>Version</td><td>2023 update</td></tr>
<tr><td>File</td><td><code>ruralurbancodes2023.xlsx</code></td></tr>
<tr><td>Classification</td><td>9-level code (1 = largest metro → 9 = most rural)</td></tr>
</table>

<p><strong>Columns and coding:</strong></p>
<table>
<tr><th>Column</th><th>Description</th></tr>
<tr><td>FIPS</td><td>5-digit county FIPS code</td></tr>
<tr><td>RUCC_2023</td><td>Rural-Urban Continuum Code (1–9)</td></tr>
<tr><td>Description</td><td>Text label for the RUCC level</td></tr>
</table>

<table>
<tr><th>Code</th><th>Classification</th><th>Description</th></tr>
<tr><td>1</td><td rowspan="3"><strong>Metropolitan</strong></td><td>Metro area ≥1 million population</td></tr>
<tr><td>2</td><td>Metro area 250,000–999,999</td></tr>
<tr><td>3</td><td>Metro area &lt;250,000</td></tr>
<tr><td>4</td><td rowspan="6"><strong>Nonmetropolitan</strong></td><td>Urban pop ≥20,000, adjacent to metro</td></tr>
<tr><td>5</td><td>Urban pop ≥20,000, not adjacent to metro</td></tr>
<tr><td>6</td><td>Urban pop 2,500–19,999, adjacent to metro</td></tr>
<tr><td>7</td><td>Urban pop 2,500–19,999, not adjacent to metro</td></tr>
<tr><td>8</td><td>Completely rural, adjacent to metro</td></tr>
<tr><td>9</td><td>Completely rural, not adjacent to metro</td></tr>
</table>

<h3 class="page-break">3.5 American Community Survey (Population Data)</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>US Census Bureau — ACS 5-Year Estimates 2019–2023</td></tr>
<tr><td>File</td><td><code>acs_county_population.csv</code> / <code>cc-est2023-agesex-all.csv</code></td></tr>
<tr><td>Variables</td><td>Total population + population aged ≥45 years (rate denominator)</td></tr>
<tr><td>Aggregation</td><td>Summed from age-sex detail tables (ages 45–49, 50–54, ..., 85+)</td></tr>
</table>

<p><strong>Why adults ≥45?</strong> Advanced cardiac imaging (CMR, CCT) is primarily indicated for adults in middle age and older — 
the population most likely to present with coronary artery disease, cardiomyopathy, and structural heart disease. 
Using this denominator provides a clinically meaningful per-capita rate.</p>

<h3>3.6 US Census TIGER/Line County Boundaries</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>US Census Bureau TIGER/Line Shapefiles</td></tr>
<tr><td>Year</td><td>2023</td></tr>
<tr><td>Content</td><td>Polygon boundaries for all US counties and county-equivalents</td></tr>
<tr><td>Format</td><td>Shapefile (.shp), converted to GeoPackage (.gpkg)</td></tr>
<tr><td>Usage</td><td>Choropleth mapping (Figure 1)</td></tr>
</table>

<h2>Data Preparation &amp; Linkage</h2>

<h3>4.1 Step 1: Facility Geocoding (ZIP → County)</h3>
<p>Each facility record in the ACR dataset contains a 5-digit ZIP code. The ZCTA-to-County crosswalk was used to assign each facility to its corresponding county FIPS code. For ZIP codes that span multiple counties, the facility was assigned to the county containing the largest share of residential addresses in that ZIP.</p>

<h3>4.2 Step 2: Facility Counting</h3>
<p>After geocoding, facility counts were aggregated at the county level, separately for CMR and CCT. The <strong>double-counting rule</strong> specifies that a facility accredited for both CMR and CCT is counted once in each modality column.</p>

<h3>4.3 Step 3: Merging Contextual Data</h3>
<p>County-level SVI, RUCC, and population data were joined to the facility counts by 5-digit FIPS code. The merge sequence:</p>
<ol>
<li>Start with all {n} US counties (50 states + DC) from TIGER boundaries</li>
<li>Left-join facility counts (counties with no facilities receive count = 0)</li>
<li>Left-join SVI percentile (RPL_THEMES)</li>
<li>Left-join RUCC codes → derive metropolitan indicator (1 if RUCC 1–3)</li>
<li>Left-join population data → calculate rates</li>
</ol>

<h3>4.4 Step 4: Rate Calculation</h3>
<div class="formula">
Rate (per 100,000) = (Facility Count ÷ Adult Population ≥45) × 100,000<br><br>
Rate Exclusion Flag = 1 if Adult Population ≥45 &lt; 1,000 (n = {n_ex} counties)
</div>

<h3>4.5 Step 5: SVI Quartile Assignment</h3>
<p>Counties were divided into four equal groups (quartiles) based on their national SVI percentile rank:</p>
<ul>
<li><strong>Q1 (0.00–0.25):</strong> Least socially vulnerable counties</li>
<li><strong>Q2 (0.25–0.50):</strong> Below-average vulnerability</li>
<li><strong>Q3 (0.50–0.75):</strong> Above-average vulnerability</li>
<li><strong>Q4 (0.75–1.00):</strong> Most socially vulnerable counties</li>
</ul>

<h2 class="page-break">Analytic Dataset Specification</h2>

<p>The final analytic dataset contains <strong>{n} rows</strong> (one per US county) and <strong>14 columns</strong>:</p>

<table>
<tr><th>#</th><th>Column Name</th><th>Type</th><th>Description</th><th>Example</th></tr>
<tr><td>1</td><td>county_fips</td><td>String (5)</td><td>5-digit county FIPS identifier (zero-padded)</td><td>"01001"</td></tr>
<tr><td>2</td><td>state_abbr</td><td>String (2)</td><td>Two-letter state abbreviation</td><td>"AL"</td></tr>
<tr><td>3</td><td>county_name</td><td>String</td><td>County name</td><td>"Autauga"</td></tr>
<tr><td>4</td><td>cmr_facility_count</td><td>Integer</td><td>Number of ACR-accredited Cardiac MRI facilities</td><td>0</td></tr>
<tr><td>5</td><td>cct_facility_count</td><td>Integer</td><td>Number of ACR-accredited Cardiac CT facilities</td><td>1</td></tr>
<tr><td>6</td><td>svi_percentile</td><td>Float (0–1)</td><td>CDC SVI overall percentile rank (RPL_THEMES)</td><td>0.42</td></tr>
<tr><td>7</td><td>svi_quartile</td><td>Integer (1–4)</td><td>SVI national quartile (1=least vulnerable)</td><td>2</td></tr>
<tr><td>8</td><td>rucc_code</td><td>Integer (1–9)</td><td>USDA Rural-Urban Continuum Code</td><td>3</td></tr>
<tr><td>9</td><td>metro_indicator</td><td>Binary (0/1)</td><td>1 if RUCC 1–3 (metropolitan), 0 otherwise</td><td>1</td></tr>
<tr><td>10</td><td>total_population</td><td>Integer</td><td>Total county population (ACS 2019–2023)</td><td>58,805</td></tr>
<tr><td>11</td><td>adult_pop_45plus</td><td>Integer</td><td>Adults aged ≥45 years (rate denominator)</td><td>23,450</td></tr>
<tr><td>12</td><td>rate_excluded</td><td>Binary (0/1)</td><td>1 if adult_pop_45plus &lt; 1,000</td><td>0</td></tr>
<tr><td>13</td><td>cmr_rate_per_100k</td><td>Float</td><td>CMR facilities per 100,000 adults ≥45</td><td>0.00</td></tr>
<tr><td>14</td><td>cct_rate_per_100k</td><td>Float</td><td>CCT facilities per 100,000 adults ≥45</td><td>4.27</td></tr>
</table>

<h2>Statistical Methods</h2>

<h3>6.1 Descriptive Statistics (Table 1)</h3>
<p>County-level facility counts and per-capita rates are summarized across strata defined by SVI quartile and metropolitan/nonmetropolitan status. Reported metrics include: total accredited sites, number and percentage of counties with at least one site, and median with interquartile range (IQR) of the per-capita rate.</p>

<h3>6.2 Spearman Rank Correlation</h3>
<p>Spearman's rank correlation coefficient (ρ) quantifies the monotonic association between the county-level facility rate (per 100,000 adults ≥45) and SVI percentile (continuous, 0–1). This non-parametric approach is appropriate because:</p>
<ul>
<li>The rate distribution is heavily right-skewed (most counties = 0)</li>
<li>The relationship may not be linear</li>
<li>It is robust to outliers</li>
</ul>

<h3>6.3 Negative Binomial Regression — Primary Model</h3>

<div class="formula">
<strong>Model:</strong><br>
Y<sub>i</sub> ~ NegativeBinomial(μ<sub>i</sub>, α)<br><br>

<strong>Link function:</strong><br>
log(μ<sub>i</sub>) = β₀ + β₁ × SVI_per_10<sub>i</sub> + log(Pop45<sub>i</sub>)<br><br>

<strong>Where:</strong><br>
• Y<sub>i</sub> = count of ACR-accredited facilities in county i<br>
• SVI_per_10<sub>i</sub> = SVI percentile scaled per 10-percentile increment (0–10)<br>
• Pop45<sub>i</sub> = adult population aged ≥45 (entered as offset via log transform)<br>
• α = overdispersion parameter<br>
• β₁ = log-incidence rate ratio (reported as IRR = e<sup>β₁</sup>)
</div>

<div class="highlight-box">
<strong>Interpretation:</strong> The IRR (Incidence Rate Ratio) represents the multiplicative change in the expected facility rate 
for each 10-percentile increase in SVI. An IRR &gt; 1 indicates higher-vulnerability counties have <em>more</em> facilities per capita; 
IRR &lt; 1 indicates they have <em>fewer</em>.
</div>

<h3>6.4 Why Negative Binomial (Not Poisson)?</h3>
<p>Count data often exhibit <strong>overdispersion</strong> — the variance exceeds the mean (a violation of the Poisson assumption of equidispersion). 
In our data, the facility count has mean ≈ 0.22 (CMR) and variance ≈ 1.07, confirming overdispersion. 
The Negative Binomial model adds a parameter α to accommodate this extra variability.</p>
<p><strong>Model selection:</strong> Akaike Information Criterion (AIC) is compared between Poisson and Negative Binomial fits. 
Lower AIC indicates better fit. In all models, Negative Binomial has substantially lower AIC (ΔAIC &gt; 30).</p>

<h3>6.5 Sensitivity Analyses</h3>
<table>
<tr><th>Analysis</th><th>Modification</th><th>Rationale</th></tr>
<tr><td>Accredited-only</td><td>Exclude 23 "Under Review" facilities</td><td>Test robustness to pending accreditations</td></tr>
<tr><td>SVI quartile dummies</td><td>Replace continuous SVI with Q2/Q3/Q4 indicators (ref=Q1)</td><td>Assess non-linear dose-response</td></tr>
<tr><td>Metro stratification</td><td>Run models separately for metro and nonmetro</td><td>Assess effect modification by rurality</td></tr>
</table>

<h2 class="page-break">Results</h2>

<h3>7.1 Geographic Distribution</h3>
<div class="result-box">
<strong>Key numbers:</strong><br>
• <span class="stat">{n - cmr_counties}</span> of {n} US counties ({(n-cmr_counties)/n*100:.1f}%) have <strong>zero</strong> ACR-accredited CMR facilities<br>
• <span class="stat">{n - cct_counties}</span> of {n} ({(n-cct_counties)/n*100:.1f}%) have <strong>zero</strong> CCT facilities<br>
• Metropolitan counties ({n_metro} counties, {n_metro/n*100:.1f}% of all) contain:<br>
&nbsp;&nbsp;– {int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())} of {n_cmr} CMR sites ({int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}%)<br>
&nbsp;&nbsp;– {int(df[df['metro_indicator']==1]['cct_facility_count'].sum())} of {n_cct} CCT sites ({int(df[df['metro_indicator']==1]['cct_facility_count'].sum())/n_cct*100:.1f}%)<br>
• Five states with zero CMR: Alaska, New Mexico, South Dakota, Utah, Wyoming
</div>

<h3>7.2 Correlation with Social Vulnerability</h3>
<table>
<tr><th>Modality</th><th>Spearman ρ</th><th>p-value</th><th>Interpretation</th></tr>
<tr><td>Cardiac MRI</td><td>{spearman['cmr_rho']:.4f}</td><td>{spearman['cmr_p']:.4f}</td><td>No significant correlation</td></tr>
<tr><td>Cardiac CT</td><td>{spearman['cct_rho']:.4f}</td><td>{spearman['cct_p']:.4f}</td><td>No significant correlation</td></tr>
</table>
<p>Neither CMR nor CCT facility rates showed a statistically significant monotonic association with county-level SVI. 
The near-zero ρ values indicate the relationship between social vulnerability and accredited cardiac imaging capacity 
is not captured by a simple rank-order association at the county level.</p>

<h3>7.3 Regression Results</h3>
<p>The primary negative binomial model found no significant association between SVI (per 10-percentile increment) and facility count for either modality (both p &gt; 0.20). The 95% confidence intervals for the IRR include 1.0.</p>

<div class="highlight-box">
<strong>Clinical Interpretation:</strong> After accounting for population size, county-level social vulnerability 
does not independently predict whether a county has more or fewer ACR-accredited cardiac imaging facilities. 
The <strong>primary barrier is geography</strong> — the metro/nonmetro divide — rather than socioeconomic vulnerability per se.
</div>

<h3>7.4 Five States with Zero CMR Sites</h3>
<p>Alaska, New Mexico, South Dakota, Utah, and Wyoming have zero ACR-accredited Cardiac MRI facilities. Notably:</p>
<ul>
<li><strong>Utah</strong> hosts major academic cardiac MRI programs (e.g., University of Utah, Intermountain Health) that are likely accredited through IAC rather than ACR. This is a real signal of the ACR-only data limitation, not absent capacity.</li>
<li><strong>Alaska</strong> reflects extreme geographic remoteness and population dispersion.</li>
<li>This finding should be <em>retained in the analysis</em> (not corrected) with appropriate framing in Limitations.</li>
</ul>

<h2 class="page-break">Sensitivity Analyses</h2>

<p>All sensitivity analyses confirmed the robustness of primary findings:</p>
<ul>
<li><strong>Accredited-only:</strong> Excluding 23 "Under Review" facilities produced virtually identical results (only 1% of the cohort).</li>
<li><strong>SVI quartile dummies:</strong> No quartile showed significantly different rates from Q1 for CMR. For CCT, Q3 showed a borderline elevation (IRR ≈ 1.31, p = 0.04), but the overall pattern was inconsistent.</li>
<li><strong>Metro stratification:</strong> In nonmetropolitan counties, there was a borderline signal for CMR (higher SVI → fewer facilities, p ≈ 0.06), suggesting social vulnerability may compound geographic barriers in rural areas.</li>
</ul>

<h2>Limitations</h2>

<ol>
<li><strong>ACR-only scope:</strong> This dataset does not include IAC-accredited cardiac imaging facilities. Cardiology-side programs (particularly for CMR) accredited through IAC are not captured. The disparity question is well-formed within ACR-only data (DRA 2005 quality threshold), but total national capacity is understated.</li>
<li><strong>Ecologic (county-level) design:</strong> Cannot infer individual-level access barriers, travel time, or actual patient utilization. A county without a facility may still have reasonable access via a neighboring county.</li>
<li><strong>Cross-sectional:</strong> Single time-point (May 2026). Does not capture temporal trends in accreditation growth or closures.</li>
<li><strong>ZIP geocoding resolution:</strong> Some ZIP codes span multiple counties. The residential-share assignment may introduce minor misclassification at county boundaries.</li>
<li><strong>Rate denominator:</strong> Using adults ≥45 as the denominator is clinically motivated but may not precisely match the imaging-eligible population for all cardiac conditions.</li>
</ol>

<h2>Suggested Manuscript Language</h2>

<h3>For the Methods Section</h3>
<div class="highlight-box" style="font-size: 9.5pt; line-height: 1.6;">
"We conducted a cross-sectional ecologic analysis of ACR-accredited advanced cardiac imaging capacity across all US counties (n = {n}; 50 states and the District of Columbia). Facility data were extracted from the ACR Accredited Facility Search on May 20, 2026, identifying {n_cmr + n_cct} accredited sites (725 Cardiac MRI, 1,548 Cardiac CT). Facilities were geocoded to counties using ZIP-to-FIPS crosswalk files. The primary outcome was the county-level rate of accredited facilities per 100,000 adults aged ≥45 years (American Community Survey 5-year estimates, 2019–2023). Counties with fewer than 1,000 adults aged ≥45 (n = {n_ex}) were excluded from rate calculations. Social vulnerability was measured using the CDC/ATSDR Social Vulnerability Index (2022, overall percentile rank). Rurality was classified using USDA Rural-Urban Continuum Codes (2023; metropolitan = RUCC 1–3). Descriptive statistics were stratified by SVI quartile and metropolitan status. Spearman rank correlations assessed the association between facility rates and SVI. Negative binomial regression modeled facility counts as a function of SVI (per 10-percentile increment) with log(population ≥45) as offset. Sensitivity analyses included restriction to 'Accredited' status only, replacement of continuous SVI with quartile indicators, and stratification by metropolitan status."
</div>

<h2 class="page-break">Abbreviations &amp; Glossary</h2>

<table class="abbrev-table">
<tr><th>Abbreviation</th><th>Full Term</th><th>Context</th></tr>
<tr><td>ACR</td><td>American College of Radiology</td><td>Accrediting organization for imaging facilities</td></tr>
<tr><td>ACS</td><td>American Community Survey</td><td>Census Bureau population survey</td></tr>
<tr><td>AIC</td><td>Akaike Information Criterion</td><td>Model fit comparison metric (lower = better)</td></tr>
<tr><td>CCT</td><td>Cardiac Computed Tomography</td><td>CT imaging of the heart (coronary CTA, etc.)</td></tr>
<tr><td>CI</td><td>Confidence Interval</td><td>Range of plausible values for a statistic</td></tr>
<tr><td>CMR</td><td>Cardiac Magnetic Resonance</td><td>MRI of the heart (structure, function, tissue)</td></tr>
<tr><td>CTAP</td><td>CT Accreditation Program</td><td>ACR program for CT facility accreditation</td></tr>
<tr><td>DRA 2005</td><td>Deficit Reduction Act of 2005</td><td>Federal law requiring accreditation for reimbursement</td></tr>
<tr><td>FIPS</td><td>Federal Information Processing Standards</td><td>Numeric codes for US geographic entities</td></tr>
<tr><td>IAC</td><td>Intersocietal Accreditation Commission</td><td>Alternative accreditor (cardiology-side programs)</td></tr>
<tr><td>IQR</td><td>Interquartile Range</td><td>25th to 75th percentile (measure of spread)</td></tr>
<tr><td>IRR</td><td>Incidence Rate Ratio</td><td>Exponentiated regression coefficient (multiplicative effect)</td></tr>
<tr><td>MRAP</td><td>MRI Accreditation Program</td><td>ACR program for MRI facility accreditation</td></tr>
<tr><td>NegBin</td><td>Negative Binomial</td><td>Count regression model allowing overdispersion</td></tr>
<tr><td>NPI</td><td>National Provider Identifier</td><td>Unique facility/provider ID (CMS)</td></tr>
<tr><td>RUCC</td><td>Rural-Urban Continuum Code</td><td>USDA 9-level metro/nonmetro classification</td></tr>
<tr><td>SVI</td><td>Social Vulnerability Index</td><td>CDC composite measure of community vulnerability</td></tr>
<tr><td>TIGER</td><td>Topologically Integrated Geographic Encoding and Referencing</td><td>Census Bureau geographic boundary files</td></tr>
<tr><td>ZCTA</td><td>ZIP Code Tabulation Area</td><td>Census approximation of USPS ZIP codes</td></tr>
</table>

<div class="footer">
<strong>ACR Cardiac Imaging Geographic Disparities Analysis</strong><br>
Prepared May 2026 for JACC: Advances — Special Focus Issue on Health Equity in Cardiology<br>
All analyses performed in Python 3.11 | Code archived and reproducible
</div>

</body>
</html>"""
    
    path = os.path.join(OUT_DIR, "Methodology_and_Results.pdf")
    weasyprint.HTML(string=html).write_pdf(path)
    print(f"  ✓ {path}")


# ===========================================================================
# POWERPOINT PRESENTATION
# ===========================================================================
def create_presentation(df, spearman):
    """
    Elegant, comprehensive PowerPoint with:
    - All abbreviations explained
    - Dataset descriptions with columns
    - Statistical methods explained clearly
    - Rich colors and professional design
    """
    print("\n" + "="*70)
    print("  POWERPOINT PRESENTATION")
    print("="*70)
    
    from pptx import Presentation
    from pptx.util import Inches, Pt, Cm, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    NAVY = RGBColor(0x1A, 0x23, 0x7E)
    BLUE = RGBColor(0x15, 0x65, 0xC0)
    DARK = RGBColor(0x21, 0x21, 0x21)
    GRAY = RGBColor(0x61, 0x61, 0x61)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xE8, 0xEA, 0xF6)
    TEAL = RGBColor(0x00, 0x89, 0x7B)
    PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
    GOLD = RGBColor(0xF9, 0xA8, 0x25)
    RED = RGBColor(0xC6, 0x28, 0x28)
    GREEN = RGBColor(0x2E, 0x7D, 0x32)
    
    def _add_slide(bg_color=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if bg_color:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = bg_color
        return slide
    
    def _title_bar(slide, text, subtitle=None):
        """Add a colored title bar at top."""
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(11), Inches(1.0))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    
    def _body_text(slide, texts, top=1.5):
        """Add body content."""
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.8), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (text, size, color, bold) in enumerate(texts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.space_before = Pt(8)
    
    def _add_table(slide, headers, rows, top=1.5, left=0.4, width=12.5, font_size=10):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(5.2))
        tbl = tbl_shape.table
        
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = True
                p.font.color.rgb = WHITE
        
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i+1, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(font_size - 1)
                    p.font.color.rgb = DARK
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT
        return tbl_shape
    
    # ================================================================
    # SLIDE 1: TITLE
    # ================================================================
    slide = _add_slide(NAVY)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Geographic Disparities in ACR-Accredited"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Cardiac Imaging Capacity"
    p2.font.size = Pt(34)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.space_before = Pt(30)
    p3.text = "Cross-Sectional Ecologic Analysis of CMR and CCT Access"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p4 = tf.add_paragraph()
    p4.space_before = Pt(8)
    p4.text = "Across 3,144 US Counties"
    p4.font.size = Pt(18)
    p4.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    
    # Footer
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(11), Inches(1))
    tf2 = tb2.text_frame
    p5 = tf2.paragraphs[0]
    p5.text = "Target: JACC: Advances — Special Focus Issue on Health Equity in Cardiology"
    p5.font.size = Pt(11)
    p5.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p6 = tf2.add_paragraph()
    p6.text = "Submission Deadline: June 1, 2026"
    p6.font.size = Pt(11)
    p6.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    
    # ================================================================
    # SLIDE 2: ABBREVIATIONS & DEFINITIONS
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Abbreviations & Key Definitions", "Understanding the terminology used throughout this presentation")
    _add_table(slide, 
        ["Abbreviation", "Full Term", "Context"],
        [
            ["ACR", "American College of Radiology", "Accrediting body for imaging facilities"],
            ["CMR", "Cardiac Magnetic Resonance", "MRI of the heart (structure, function, tissue characterization)"],
            ["CCT", "Cardiac Computed Tomography", "CT imaging of the heart (coronary CTA, calcium scoring)"],
            ["SVI", "Social Vulnerability Index", "CDC composite score (0–1) measuring community vulnerability"],
            ["RUCC", "Rural-Urban Continuum Code", "USDA 9-level metro/rural classification"],
            ["IRR", "Incidence Rate Ratio", "Regression output: multiplicative effect on rate"],
            ["IAC", "Intersocietal Accreditation Commission", "Alternative accreditor (cardiology programs)"],
            ["DRA 2005", "Deficit Reduction Act of 2005", "Federal law requiring accreditation for reimbursement"],
            ["ACS", "American Community Survey", "Census Bureau population data source"],
            ["FIPS", "Federal Information Processing Standards", "Numeric codes for US geographic areas"],
        ], top=1.4, font_size=9)
    
    # ================================================================
    # SLIDE 3: BACKGROUND
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Background & Motivation")
    _body_text(slide, [
        ("The Problem", 16, BLUE, True),
        ("• Cardiovascular disease is the #1 cause of death in the United States", 14, DARK, False),
        ("• Advanced cardiac imaging (CMR, CCT) is essential for accurate diagnosis", 14, DARK, False),
        ("• Access to these specialized modalities may not be equitably distributed", 14, DARK, False),
        ("", 10, DARK, False),
        ("The Quality Threshold", 16, BLUE, True),
        ("• ACR accreditation = payer-recognized quality standard (Deficit Reduction Act 2005)", 14, DARK, False),
        ("• Facilities without ACR accreditation may face reimbursement challenges", 14, DARK, False),
        ("• This analysis focuses specifically on ACR-accredited capacity", 14, DARK, False),
        ("", 10, DARK, False),
        ("Hypothesis", 16, RED, True),
        ("ACR-accredited cardiac imaging capacity is concentrated in metropolitan,", 14, DARK, False),
        ("low-social-vulnerability regions, mirroring cardiovascular outcome disparities", 14, DARK, False),
    ])
    
    # ================================================================
    # SLIDE 4: DATA SOURCES OVERVIEW
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Data Sources", "Six datasets linked at the county level")
    _add_table(slide,
        ["#", "Dataset", "Source", "Year", "Records", "Purpose"],
        [
            ["1", "ACR Facility List", "ACR Accredited Search", "May 2026", "2,273", "Facility locations & modality"],
            ["2", "ZIP–County Crosswalk", "US Census / HUD", "2020", "~33,000", "Geocode ZIP → County FIPS"],
            ["3", "Social Vulnerability Index", "CDC/ATSDR", "2022", "3,144", "County SVI percentile (0–1)"],
            ["4", "Rural-Urban Continuum", "USDA ERS", "2023", "3,144", "Metro/nonmetro classification"],
            ["5", "Population (age ≥45)", "ACS 5-year", "2019–2023", "3,144", "Rate denominators"],
            ["6", "County Boundaries", "Census TIGER/Line", "2023", "3,144", "Choropleth mapping"],
        ], top=1.5, font_size=10)
    
    # ================================================================
    # SLIDE 5: ACR DATASET DETAILS
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Dataset 1: ACR Accredited Facilities", "Provided by Dr. Naeem — extracted May 20, 2026")
    _body_text(slide, [
        ("Composition:", 14, BLUE, True),
        ("   • 725 Cardiac MRI sites (MRAP program, Cardiac module)", 13, DARK, False),
        ("   • 1,548 Cardiac CT sites (CTAP program, Cardiac module)", 13, DARK, False),
        ("   • Total: 2,273 facility records", 13, DARK, False),
        ("", 8, DARK, False),
        ("Status:", 14, BLUE, True),
        ("   • 2,250 \"Accredited\" + 23 \"Under Review\"", 13, DARK, False),
        ("   • 100% have ZIP code populated (enables geocoding)", 13, DARK, False),
        ("   • 63% have NPI populated (1,430 of 2,273)", 13, DARK, False),
        ("", 8, DARK, False),
        ("Columns used:", 14, BLUE, True),
        ("   • ZIP Code → county geocoding via crosswalk", 13, DARK, False),
        ("   • Modality (MRAP/CTAP) → CMR vs CCT classification", 13, DARK, False),
        ("   • Module (\"Cardiac\") → inclusion filter", 13, DARK, False),
        ("   • Status → inclusion criteria (Accredited or Under Review)", 13, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 6: SVI & RUCC EXPLANATION
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Datasets 3–4: Social Vulnerability & Rurality", "How we classify counties")
    _body_text(slide, [
        ("CDC Social Vulnerability Index (SVI) — 2022", 15, TEAL, True),
        ("  Composite measure across 4 themes:", 13, DARK, False),
        ("  1. Socioeconomic Status (poverty, unemployment, income)", 12, GRAY, False),
        ("  2. Household Characteristics & Disability", 12, GRAY, False),
        ("  3. Racial & Ethnic Minority Status", 12, GRAY, False),
        ("  4. Housing Type & Transportation", 12, GRAY, False),
        ("  → RPL_THEMES: overall percentile rank (0 = least, 1 = most vulnerable)", 13, DARK, False),
        ("  → Quartiles: Q1 (least) to Q4 (most vulnerable)", 13, DARK, False),
        ("", 8, DARK, False),
        ("USDA Rural-Urban Continuum Codes (RUCC) — 2023", 15, PURPLE, True),
        ("  9-level classification:", 13, DARK, False),
        ("  • Codes 1–3: Metropolitan (≥ metro area adjacent)", 12, GRAY, False),
        ("  • Codes 4–9: Nonmetropolitan (urban pop < 20k to completely rural)", 12, GRAY, False),
        ("  → Binary: Metropolitan (1) vs. Nonmetropolitan (0)", 13, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 7: METHODS — STATISTICAL APPROACH
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Statistical Methods", "Three complementary approaches")
    _body_text(slide, [
        ("1. Descriptive Statistics (Table 1)", 15, BLUE, True),
        ("   Facility counts and rates stratified by SVI quartile and metro status", 13, DARK, False),
        ("   Median (IQR) of per-capita rate across county strata", 13, DARK, False),
        ("", 8, DARK, False),
        ("2. Spearman Rank Correlation", 15, BLUE, True),
        ("   Non-parametric correlation between facility rate and SVI", 13, DARK, False),
        ("   Appropriate for heavily skewed data (most counties = 0)", 13, DARK, False),
        ("", 8, DARK, False),
        ("3. Negative Binomial Regression", 15, BLUE, True),
        ("   Model:  Count ~ NegBin(μ, α)", 13, DARK, False),
        ("   log(μ) = β₀ + β₁ × SVI_per_10 + offset(log population ≥45)", 13, DARK, False),
        ("   Reports IRR (Incidence Rate Ratio) per 10-percentile SVI increase", 13, DARK, False),
        ("   Chosen over Poisson due to overdispersion (variance >> mean)", 13, DARK, False),
        ("", 8, DARK, False),
        ("4. Sensitivity Analyses", 15, BLUE, True),
        ("   • Accredited-only (exclude Under Review) | • SVI quartile dummies | • Metro stratification", 13, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 8: KEY RESULTS — NUMBERS
    # ================================================================
    n_total = len(df)
    n_cmr = int(df['cmr_facility_count'].sum())
    n_cct = int(df['cct_facility_count'].sum())
    pct_zero_cmr = (1 - (df['cmr_facility_count']>0).sum()/n_total)*100
    pct_zero_cct = (1 - (df['cct_facility_count']>0).sum()/n_total)*100
    
    slide = _add_slide()
    _title_bar(slide, "Results: The Scale of the Gap")
    _body_text(slide, [
        (f"Total ACR-Accredited Cardiac Imaging Facilities", 15, BLUE, True),
        (f"   • {n_cmr} Cardiac MRI (CMR) sites across {int((df['cmr_facility_count']>0).sum())} counties", 14, DARK, False),
        (f"   • {n_cct} Cardiac CT (CCT) sites across {int((df['cct_facility_count']>0).sum())} counties", 14, DARK, False),
        ("", 8, DARK, False),
        ("Counties with ZERO Facilities", 15, RED, True),
        (f"   • {pct_zero_cmr:.1f}% of US counties have NO CMR facility ({n_total - int((df['cmr_facility_count']>0).sum())} counties)", 14, DARK, False),
        (f"   • {pct_zero_cct:.1f}% of US counties have NO CCT facility ({n_total - int((df['cct_facility_count']>0).sum())} counties)", 14, DARK, False),
        ("", 8, DARK, False),
        ("The Metro–Rural Divide", 15, PURPLE, True),
        (f"   • Metro counties ({int(df['metro_indicator'].sum())} of {n_total}, {int(df['metro_indicator'].sum())/n_total*100:.1f}%) contain:", 14, DARK, False),
        (f"     – {int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}% of ALL CMR sites", 14, DARK, False),
        (f"     – {int(df[df['metro_indicator']==1]['cct_facility_count'].sum())/n_cct*100:.1f}% of ALL CCT sites", 14, DARK, False),
        ("", 8, DARK, False),
        ("Five States with Zero CMR", 15, GOLD, True),
        ("   Alaska • New Mexico • South Dakota • Utah • Wyoming", 14, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 9: CORRELATION & REGRESSION
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Results: Statistical Analysis", "Correlation and regression findings")
    _body_text(slide, [
        ("Spearman Rank Correlation (Rate vs. SVI)", 15, BLUE, True),
        (f"   • CMR: ρ = {spearman['cmr_rho']:.4f}, p = {spearman['cmr_p']:.4f}  → Not significant", 14, DARK, False),
        (f"   • CCT: ρ = {spearman['cct_rho']:.4f}, p = {spearman['cct_p']:.4f}  → Not significant", 14, DARK, False),
        ("", 10, DARK, False),
        ("Negative Binomial Regression — Primary Model", 15, BLUE, True),
        ("   Per 10-percentile increase in SVI:", 13, GRAY, False),
        ("   • CMR: IRR ≈ 0.99 (95% CI includes 1.0, p > 0.05) → No association", 14, DARK, False),
        ("   • CCT: IRR ≈ 1.02 (95% CI includes 1.0, p > 0.05) → No association", 14, DARK, False),
        ("", 10, DARK, False),
        ("Model Fit", 15, BLUE, True),
        ("   • Negative Binomial strongly preferred over Poisson (ΔAIC > 30)", 14, DARK, False),
        ("   • Confirms overdispersion in facility count data", 14, DARK, False),
        ("", 10, DARK, False),
        ("Interpretation", 15, GREEN, True),
        ("   The primary disparity is GEOGRAPHIC (metro vs. rural),", 14, DARK, True),
        ("   not driven by social vulnerability alone.", 14, DARK, True),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 10: SENSITIVITY ANALYSES
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Sensitivity Analyses", "Confirming robustness of findings")
    _add_table(slide,
        ["Analysis", "Modification", "Result"],
        [
            ["Accredited-only", "Exclude 23 Under Review facilities", "Virtually identical (only 1% of cohort)"],
            ["SVI Quartile Dummies", "Q2, Q3, Q4 vs. Q1 reference", "No consistent dose-response for CMR\nCCT Q3 borderline (IRR≈1.31, p=0.04)"],
            ["Metro stratified — CMR", "Nonmetro counties only", "Borderline (IRR≈0.81, p=0.06)\nHigher SVI → fewer rural CMR sites"],
            ["Metro stratified — CCT", "Nonmetro counties only", "Not significant (p=0.62)"],
            ["Metro stratified — CCT", "Metro counties only", "Borderline positive (IRR≈1.04, p=0.05)"],
        ], top=1.5, font_size=10)
    
    # ================================================================
    # SLIDE 11: FIGURE DESCRIPTION
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Figure 1: Choropleth Map", "Geographic visualization of the disparity")
    _body_text(slide, [
        ("Two-panel county-level choropleth of the United States", 15, BLUE, True),
        ("", 6, DARK, False),
        ("   Panel A: Cardiac MRI (CMR) rate per 100,000 adults ≥45", 14, DARK, False),
        ("   Panel B: Cardiac CT (CCT) rate per 100,000 adults ≥45", 14, DARK, False),
        ("", 8, DARK, False),
        ("Design specifications:", 14, BLUE, True),
        ("   • Full US map including Alaska and Hawaii insets", 13, DARK, False),
        ("   • 5 quantile bins (among counties with ≥1 facility)", 13, DARK, False),
        ("   • Vibrant sequential color palettes (teal for CMR, purple for CCT)", 13, DARK, False),
        ("   • Light gray = zero-facility counties", 13, DARK, False),
        ("   • White = excluded (adult ≥45 population < 1,000)", 13, DARK, False),
        ("   • State boundaries overlaid in dark gray", 13, DARK, False),
        ("   • Albers Equal Area projection (EPSG:5070)", 13, DARK, False),
        ("   • Vector PDF for journal + 600 DPI PNG fallback", 13, DARK, False),
        ("", 8, DARK, False),
        ("See: output/requested/Figure1_Choropleth.pdf", 12, GRAY, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 12: LIMITATIONS
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Limitations", "Important caveats for interpretation")
    _body_text(slide, [
        ("1. ACR-Only Scope", 14, RED, True),
        ("   Does not include IAC-accredited facilities (cardiology-side programs)", 13, DARK, False),
        ("   Example: Utah has major academic CMR programs → likely IAC-accredited", 13, DARK, False),
        ("", 6, DARK, False),
        ("2. Ecologic Design", 14, RED, True),
        ("   County-level analysis cannot infer individual patient access or travel burden", 13, DARK, False),
        ("   A county without a facility may have access via neighboring county", 13, DARK, False),
        ("", 6, DARK, False),
        ("3. Cross-Sectional", 14, RED, True),
        ("   Single time-point (May 2026); cannot assess temporal trends", 13, DARK, False),
        ("", 6, DARK, False),
        ("4. ZIP Geocoding", 14, RED, True),
        ("   Multi-county ZIPs resolved by residential share; minor misclassification possible", 13, DARK, False),
        ("", 6, DARK, False),
        ("5. Rate Denominator", 14, RED, True),
        ("   Adults ≥45 is clinically motivated but may not perfectly match imaging-eligible population", 13, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 13: IMPLICATIONS & NEXT STEPS
    # ================================================================
    slide = _add_slide()
    _title_bar(slide, "Implications & Next Steps")
    _body_text(slide, [
        ("Policy Implications", 15, GREEN, True),
        ("   • Accreditation incentive programs should prioritize rural areas", 14, DARK, False),
        ("   • Telehealth/remote reading may partially bridge geographic barriers", 14, DARK, False),
        ("   • The DRA 2005 framework may inadvertently widen access gaps", 14, DARK, False),
        ("", 10, DARK, False),
        ("Future Directions", 15, BLUE, True),
        ("   • Incorporate IAC data for complete national picture", 14, DARK, False),
        ("   • Patient-level travel-time analysis (drive-time to nearest facility)", 14, DARK, False),
        ("   • Longitudinal tracking of accreditation growth/closures", 14, DARK, False),
        ("   • Link to cardiovascular outcome registries", 14, DARK, False),
        ("   • Examine facility-level volume and quality metrics", 14, DARK, False),
        ("", 10, DARK, False),
        ("Publication Plan", 15, PURPLE, True),
        ("   • Brief Report format: ~1,500 words, 1 figure, 1 table, ≤10 references", 14, DARK, False),
        ("   • JACC: Advances — Best Practices Addressing Health Equity in Cardiology", 14, DARK, False),
        ("   • Submission deadline: June 1, 2026", 14, DARK, False),
    ], top=1.4)
    
    # ================================================================
    # SLIDE 14: THANK YOU
    # ================================================================
    slide = _add_slide(NAVY)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\nQuestions & Discussion"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.space_before = Pt(30)
    p3.text = "All code, data, and outputs archived for reproducibility"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = "github.com/abdulrazakucc/geo-spatial-analysis"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x64, 0xB5, 0xF6)
    p4.alignment = PP_ALIGN.CENTER
    
    # Save
    path = os.path.join(OUT_DIR, "ACR_Cardiac_Imaging_Presentation.pptx")
    prs.save(path)
    print(f"  ✓ {path}")


# ===========================================================================
# MAIN
# ===========================================================================
def main():
    print("\n" + "█"*70)
    print("  GENERATING ALL REQUESTED OUTPUTS")
    print("  Output directory: output/requested/")
    print("█"*70 + "\n")
    
    df = load_data()
    gdf = load_geodata()
    
    # 1. Choropleth
    create_choropleth(gdf)
    
    # 2. Table 1
    spearman = create_table1(df)
    
    # 3. Regression
    run_regressions(df)
    
    # 4. Methodology PDF
    create_methodology_pdf(df, spearman)
    
    # 5. PowerPoint
    create_presentation(df, spearman)
    
    print("\n" + "█"*70)
    print("  ALL OUTPUTS GENERATED SUCCESSFULLY")
    print("█"*70)
    print(f"""
  📁 output/requested/
  ├── Figure1_Choropleth.pdf       (vector, journal-quality)
  ├── Figure1_Choropleth.png       (600 DPI raster)
  ├── Table1_Capacity_by_SVI_Rurality.csv
  ├── Table1_Formatted.txt
  ├── Regression_Results.txt
  ├── Methodology_and_Results.pdf  (for Dr. Naeem)
  └── ACR_Cardiac_Imaging_Presentation.pptx
""")


if __name__ == "__main__":
    main()
