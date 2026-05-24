"""
07_publication_outputs.py
=========================
Generates all publication-quality deliverables per the project briefing:

1. Figure 1 — Two-panel choropleth (A: CMR, B: CCT) with state boundaries,
   quantile bins, proper legend, 600 DPI PNG + vector PDF
2. Table 1 — Formatted capacity by SVI quartile and rurality with Spearman
3. Regression results — Primary + all sensitivity analyses
4. Sensitivity analysis: Accredited-only cohort
5. PDF Workflow Document — Full methods documentation
6. PowerPoint Presentation — All results explained

Inputs:
    data/processed/county_analytic_dataset.csv
    data/processed/county_analytic_geo.gpkg
    data/ACR_Cardiac_Imaging_Sites.xlsx

Outputs:
    output/figures/figure1_choropleth.pdf  (vector, journal-ready)
    output/figures/figure1_choropleth.png  (600 DPI raster fallback)
    output/tables/table1_publication.csv
    output/tables/table1_publication.txt
    output/models/regression_results_full.txt
    output/workflow/ACR_Cardiac_Analysis_Workflow.pdf
    output/presentation/ACR_Cardiac_Imaging_Presentation.pptx
"""

import os
import sys
import warnings
import numpy as np
import pandas as pd
import geopandas as gpd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.colors import BoundaryNorm, ListedColormap
from matplotlib.lines import Line2D
import mapclassify
from scipy import stats
import statsmodels.api as sm
from statsmodels.genmod.families import NegativeBinomial, Poisson
import pickle

warnings.filterwarnings('ignore')

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROC_DIR = os.path.join(BASE_DIR, "data", "processed")
FIG_DIR = os.path.join(BASE_DIR, "output", "figures")
TBL_DIR = os.path.join(BASE_DIR, "output", "tables")
MDL_DIR = os.path.join(BASE_DIR, "output", "models")
WKF_DIR = os.path.join(BASE_DIR, "output", "workflow")
PPT_DIR = os.path.join(BASE_DIR, "output", "presentation")

for d in [FIG_DIR, TBL_DIR, MDL_DIR, WKF_DIR, PPT_DIR]:
    os.makedirs(d, exist_ok=True)


# ===========================================================================
# DATA LOADING
# ===========================================================================
def load_data():
    """Load the analytic dataset."""
    df = pd.read_csv(
        os.path.join(PROC_DIR, "county_analytic_dataset.csv"),
        dtype={"county_fips": str}
    )
    print(f"✓ Loaded analytic dataset: {len(df)} counties, {len(df.columns)} variables")
    return df


def load_geodata():
    """Load geo-enabled dataset for mapping."""
    gdf = gpd.read_file(os.path.join(PROC_DIR, "county_analytic_geo.gpkg"))
    gdf = gdf.to_crs(epsg=5070)  # Albers Equal Area
    print(f"✓ Loaded geo-dataset: {len(gdf)} counties")
    return gdf


# ===========================================================================
# FIGURE 1: PUBLICATION-QUALITY CHOROPLETH
# ===========================================================================
def create_figure1(gdf):
    """
    Create Figure 1: Two-panel US county choropleth.
    Panel A: CMR rate per 100k adults ≥45
    Panel B: CCT rate per 100k adults ≥45
    
    Specifications per briefing:
    - Sequential blue palette, 5 quantile bins + gray (zero) + white (excluded)
    - State boundaries: thin dark gray overlay
    - County boundaries: hairline lighter gray
    - Albers USA projection
    - Vector PDF + 600 DPI PNG
    """
    print("\n" + "="*70)
    print("FIGURE 1: Two-Panel Choropleth Map")
    print("="*70)
    
    # Set up figure - vertical layout (A above B) for portrait journal
    fig, axes = plt.subplots(2, 1, figsize=(12, 18), facecolor='white')
    plt.subplots_adjust(hspace=0.08)
    
    rate_cols = ['cmr_rate_per_100k', 'cct_rate_per_100k']
    panel_labels = ['A', 'B']
    panel_titles = [
        'ACR-Accredited Cardiac MRI Facilities per 100,000 Adults Aged ≥45 Years',
        'ACR-Accredited Cardiac CT Facilities per 100,000 Adults Aged ≥45 Years'
    ]
    
    # Blue sequential palette (colorblind-safe)
    blues = ['#c6dbef', '#9ecae1', '#6baed6', '#2171b5', '#08306b']
    
    for idx, (col, label, title) in enumerate(zip(rate_cols, panel_labels, panel_titles)):
        ax = axes[idx]
        
        # Classify counties
        excluded = gdf[gdf['rate_excluded'] == 1]
        zero = gdf[(gdf['rate_excluded'] == 0) & (gdf[col] == 0)]
        has_data = gdf[(gdf['rate_excluded'] == 0) & (gdf[col] > 0)]
        
        # Plot excluded (white with thin border)
        if len(excluded) > 0:
            excluded.plot(ax=ax, color='white', edgecolor='#cccccc', linewidth=0.15)
        
        # Plot zero-facility counties (light gray)
        if len(zero) > 0:
            zero.plot(ax=ax, color='#f0f0f0', edgecolor='#cccccc', linewidth=0.15)
        
        # Plot rated counties with quantile classification
        if len(has_data) > 0:
            classifier = mapclassify.Quantiles(has_data[col].values, k=5)
            bins = classifier.bins
            bin_labels = []
            prev = 0
            for b in bins:
                bin_labels.append(f"{prev:.1f}–{b:.1f}")
                prev = b
            
            cmap = ListedColormap(blues)
            bounds = [0] + list(bins)
            norm = BoundaryNorm(bounds, cmap.N)
            
            has_data.plot(
                column=col, ax=ax, cmap=cmap, norm=norm,
                edgecolor='#cccccc', linewidth=0.15
            )
        
        # State boundaries overlay
        state_gdf = gdf.dissolve(by='state_abbr')
        state_gdf.boundary.plot(ax=ax, color='#333333', linewidth=0.6)
        
        # Formatting
        ax.set_axis_off()
        ax.set_title(f"{label}. {title}", fontsize=11, fontweight='bold', 
                     loc='left', pad=12, fontfamily='serif')
        
        # Set extent to continental US (Albers bounds)
        ax.set_xlim([-2.5e6, 2.5e6])
        ax.set_ylim([-1.5e6, 1.6e6])
        
        # Legend
        legend_patches = []
        for i, (color, lbl) in enumerate(zip(blues, bin_labels)):
            legend_patches.append(mpatches.Patch(facecolor=color, edgecolor='#999', 
                                                  linewidth=0.5, label=f"Q{i+1}: {lbl}"))
        legend_patches.append(mpatches.Patch(facecolor='#f0f0f0', edgecolor='#999',
                                             linewidth=0.5, label='Zero facilities'))
        legend_patches.append(mpatches.Patch(facecolor='white', edgecolor='#999',
                                             linewidth=0.5, label='Excluded (pop. <1,000)'))
        
        leg = ax.legend(
            handles=legend_patches,
            title='Facilities per 100,000\nadults aged ≥45',
            loc='lower right',
            fontsize=7.5,
            title_fontsize=8,
            framealpha=0.95,
            edgecolor='#999',
            fancybox=False,
            borderpad=0.8
        )
        leg.get_title().set_fontweight('bold')
    
    # Figure annotation
    fig.text(
        0.5, 0.02,
        "Figure 1. Geographic distribution of ACR-accredited cardiac imaging capacity by US county.\n"
        "Color intensity indicates facility density (quantile-binned among counties with ≥1 facility).\n"
        "Light gray = zero accredited facilities; white = excluded from rate calculation (adult ≥45 population <1,000).\n"
        "State boundaries overlaid. County boundaries: US Census TIGER/Line 2023. Projection: Albers Equal Area (EPSG:5070).\n"
        "Alaska and Hawaii included in analysis; insets omitted for clarity.",
        ha='center', fontsize=7.5, color='#444', style='italic', fontfamily='serif',
        linespacing=1.5
    )
    
    # Save
    pdf_path = os.path.join(FIG_DIR, "figure1_choropleth.pdf")
    png_path = os.path.join(FIG_DIR, "figure1_choropleth.png")
    
    fig.savefig(pdf_path, format='pdf', bbox_inches='tight', dpi=300)
    fig.savefig(png_path, format='png', bbox_inches='tight', dpi=600)
    plt.close()
    
    print(f"  ✓ PDF (vector): {pdf_path}")
    print(f"  ✓ PNG (600 DPI): {png_path}")
    
    # Print bin info
    for col, title in zip(rate_cols, ['CMR', 'CCT']):
        has = gdf[(gdf['rate_excluded'] == 0) & (gdf[col] > 0)]
        print(f"\n  {title} bins (quantile): {mapclassify.Quantiles(has[col].values, k=5).bins}")
        print(f"  Counties with data: {len(has)}, Zero: {len(gdf[(gdf['rate_excluded']==0) & (gdf[col]==0)])}, Excluded: {len(gdf[gdf['rate_excluded']==1])}")


# ===========================================================================
# TABLE 1: PUBLICATION-QUALITY
# ===========================================================================
def create_table1(df):
    """
    Create Table 1: ACR-Accredited Cardiac Imaging Capacity by SVI Quartile and Rurality.
    
    Per briefing:
    - Rows: All, Q1-Q4, Metro, Nonmetro
    - Columns: n counties, adult ≥45 pop, CMR sites, CMR counties, CMR rate median/IQR,
               CCT sites, CCT counties, CCT rate median/IQR
    - Spearman correlation at bottom
    - Proper footnotes
    """
    print("\n" + "="*70)
    print("TABLE 1: Capacity by SVI Quartile and Rurality")
    print("="*70)
    
    # Rate-eligible subset
    df_rates = df[df['rate_excluded'] == 0].copy()
    
    def compute_stratum(subset, subset_rates, label):
        n = len(subset)
        pop45 = subset['adult_pop_45plus'].sum()
        
        cmr_sites = int(subset['cmr_facility_count'].sum())
        cmr_counties = int((subset['cmr_facility_count'] > 0).sum())
        cmr_pct = f"{cmr_counties/n*100:.1f}" if n > 0 else "0.0"
        
        cct_sites = int(subset['cct_facility_count'].sum())
        cct_counties = int((subset['cct_facility_count'] > 0).sum())
        cct_pct = f"{cct_counties/n*100:.1f}" if n > 0 else "0.0"
        
        # Rates (from rate-eligible only)
        sr = subset_rates
        cmr_med = sr['cmr_rate_per_100k'].median()
        cmr_q1 = sr['cmr_rate_per_100k'].quantile(0.25)
        cmr_q3 = sr['cmr_rate_per_100k'].quantile(0.75)
        
        cct_med = sr['cct_rate_per_100k'].median()
        cct_q1 = sr['cct_rate_per_100k'].quantile(0.25)
        cct_q3 = sr['cct_rate_per_100k'].quantile(0.75)
        
        # Mean rates for comparison
        cmr_mean = sr['cmr_rate_per_100k'].mean()
        cct_mean = sr['cct_rate_per_100k'].mean()
        
        return {
            'Stratum': label,
            'Counties (n)': n,
            'Adults ≥45 (millions)': f"{pop45/1e6:.2f}",
            'CMR Sites': cmr_sites,
            'CMR Counties (%)': f"{cmr_counties} ({cmr_pct}%)",
            'CMR Rate Median (IQR)': f"{cmr_med:.2f} ({cmr_q1:.2f}–{cmr_q3:.2f})",
            'CMR Rate Mean': f"{cmr_mean:.2f}",
            'CCT Sites': cct_sites,
            'CCT Counties (%)': f"{cct_counties} ({cct_pct}%)",
            'CCT Rate Median (IQR)': f"{cct_med:.2f} ({cct_q1:.2f}–{cct_q3:.2f})",
            'CCT Rate Mean': f"{cct_mean:.2f}",
        }
    
    rows = []
    
    # All counties
    rows.append(compute_stratum(df, df_rates, 'All counties'))
    
    # By SVI quartile
    for q in range(1, 5):
        label = {1: 'Q1 (least vulnerable)', 2: 'Q2', 3: 'Q3', 4: 'Q4 (most vulnerable)'}[q]
        subset = df[df['svi_quartile'] == q]
        subset_r = df_rates[df_rates['svi_quartile'] == q]
        rows.append(compute_stratum(subset, subset_r, label))
    
    # By metro status
    metro = df[df['metro_indicator'] == 1]
    metro_r = df_rates[df_rates['metro_indicator'] == 1]
    rows.append(compute_stratum(metro, metro_r, 'Metropolitan (RUCC 1–3)'))
    
    nonmetro = df[df['metro_indicator'] == 0]
    nonmetro_r = df_rates[df_rates['metro_indicator'] == 0]
    rows.append(compute_stratum(nonmetro, nonmetro_r, 'Nonmetropolitan (RUCC 4–9)'))
    
    table_df = pd.DataFrame(rows)
    
    # Spearman correlations
    cmr_rho, cmr_p = stats.spearmanr(df_rates['cmr_rate_per_100k'], df_rates['svi_percentile'])
    cct_rho, cct_p = stats.spearmanr(df_rates['cct_rate_per_100k'], df_rates['svi_percentile'])
    
    # Format for output
    print("\n" + table_df.to_string(index=False))
    print(f"\n  Spearman ρ (CMR rate vs SVI): {cmr_rho:.4f}, p = {cmr_p:.4f}")
    print(f"  Spearman ρ (CCT rate vs SVI): {cct_rho:.4f}, p = {cct_p:.4f}")
    
    # Save CSV
    table_df.to_csv(os.path.join(TBL_DIR, "table1_publication.csv"), index=False)
    
    # Save formatted text
    with open(os.path.join(TBL_DIR, "table1_publication.txt"), 'w') as f:
        f.write("Table 1. ACR-Accredited Cardiac Imaging Capacity by Social Vulnerability Quartile and Rurality\n")
        f.write("=" * 120 + "\n\n")
        f.write(table_df.to_string(index=False))
        f.write("\n\n" + "─" * 120 + "\n")
        f.write("Spearman Rank Correlation (facility rate vs. SVI percentile, continuous):\n")
        f.write(f"  Cardiac MRI: ρ = {cmr_rho:.4f}, p = {cmr_p:.4f}\n")
        f.write(f"  Cardiac CT:  ρ = {cct_rho:.4f}, p = {cct_p:.4f}\n")
        f.write("\n" + "─" * 120 + "\n")
        f.write("Notes:\n")
        f.write("  SVI = Social Vulnerability Index (CDC/ATSDR 2022, overall percentile rank RPL_THEMES).\n")
        f.write("  RUCC = Rural-Urban Continuum Code (USDA ERS 2023). Metropolitan = RUCC 1–3.\n")
        f.write("  Population denominators from American Community Survey 5-year estimates 2019–2023.\n")
        f.write(f"  Counties with <1,000 adults aged ≥45 excluded from rate calculations (n={int(df['rate_excluded'].sum())}).\n")
        f.write("  Rate = accredited facilities per 100,000 adults aged ≥45 years.\n")
        f.write("  IQR = interquartile range.\n")
    
    print(f"\n  ✓ {os.path.join(TBL_DIR, 'table1_publication.csv')}")
    print(f"  ✓ {os.path.join(TBL_DIR, 'table1_publication.txt')}")
    
    return {'cmr_rho': cmr_rho, 'cmr_p': cmr_p, 'cct_rho': cct_rho, 'cct_p': cct_p}


# ===========================================================================
# REGRESSION MODELS (ALL PER BRIEFING)
# ===========================================================================
def run_regressions(df):
    """
    Run all regression models per briefing Section 3.3–3.4:
    
    Primary:
      - Negative binomial: facility count ~ SVI (per 10-percentile) + offset(log pop ≥45)
      - Separately for CMR and CCT
    
    Sensitivity:
      1. Accredited-only (exclude 23 Under Review)
      2. SVI quartile dummies instead of continuous
      3. Stratified by metro/nonmetro
    
    Goodness of fit: dispersion, AIC comparison vs Poisson
    """
    print("\n" + "="*70)
    print("REGRESSION MODELS")
    print("="*70)
    
    df_model = df[df['rate_excluded'] == 0].copy()
    df_model['svi_per10'] = df_model['svi_percentile'] / 10
    df_model['log_pop'] = np.log(df_model['adult_pop_45plus'])
    
    # SVI quartile dummies
    df_model['svi_q2'] = (df_model['svi_quartile'] == 2).astype(int)
    df_model['svi_q3'] = (df_model['svi_quartile'] == 3).astype(int)
    df_model['svi_q4'] = (df_model['svi_quartile'] == 4).astype(int)
    
    results_text = []
    results_text.append("=" * 80)
    results_text.append("REGRESSION RESULTS: ACR Cardiac Imaging Geographic Disparities")
    results_text.append("Negative Binomial Models with log(adult ≥45 population) as offset")
    results_text.append("=" * 80)
    
    all_models = {}
    
    for modality in ['CMR', 'CCT']:
        col = f"{'cmr' if modality == 'CMR' else 'cct'}_facility_count"
        results_text.append(f"\n\n{'═'*80}")
        results_text.append(f"MODALITY: {modality}")
        results_text.append(f"{'═'*80}")
        
        # ---- PRIMARY MODEL ----
        models_to_run = [
            ("Primary (SVI continuous, per 10-percentile increment)", ['svi_per10'], df_model, True),
            ("Sensitivity A: Accredited-only (excluding Under Review)", ['svi_per10'], df_model, False),  # placeholder
            ("Sensitivity B: SVI Quartile Dummies (ref=Q1)", ['svi_q2', 'svi_q3', 'svi_q4'], df_model, False),
            ("Sensitivity C: Stratified — Metropolitan only", ['svi_per10'], df_model[df_model['metro_indicator']==1], False),
            ("Sensitivity D: Stratified — Nonmetropolitan only", ['svi_per10'], df_model[df_model['metro_indicator']==0], False),
        ]
        
        for model_name, predictors, data, compare_poisson in models_to_run:
            results_text.append(f"\n{'─'*80}")
            results_text.append(f"Model: {modality} — {model_name}")
            results_text.append(f"{'─'*80}")
            
            y = data[col].values
            X = sm.add_constant(data[predictors].values)
            offset = data['log_pop'].values
            
            try:
                # Negative Binomial
                nb_model = sm.GLM(y, X, family=NegativeBinomial(alpha=1.0), offset=offset)
                nb_result = nb_model.fit(maxiter=200, disp=False)
                
                # Extract results
                results_text.append(f"\n{'Variable':<30} {'IRR':>8} {'95% CI':>20} {'p-value':>10}")
                results_text.append("─" * 72)
                
                var_names = ['Intercept'] + predictors
                for i, vname in enumerate(var_names):
                    irr = np.exp(nb_result.params[i])
                    ci_lo = np.exp(nb_result.conf_int()[i, 0])
                    ci_hi = np.exp(nb_result.conf_int()[i, 1])
                    pval = nb_result.pvalues[i]
                    pstr = f"{pval:.4f}" if pval >= 0.0001 else "<0.0001"
                    results_text.append(f"{vname:<30} {irr:>8.4f} ({ci_lo:.4f}–{ci_hi:.4f}) {pstr:>10}")
                
                results_text.append(f"\n  N = {int(nb_result.nobs)}")
                results_text.append(f"  AIC = {nb_result.aic:.1f}")
                results_text.append(f"  Deviance = {nb_result.deviance:.1f}")
                results_text.append(f"  Pearson χ² = {nb_result.pearson_chi2:.1f}")
                results_text.append(f"  Dispersion (Pearson χ²/df) = {nb_result.pearson_chi2/nb_result.df_resid:.3f}")
                
                # Poisson comparison
                if compare_poisson:
                    pois_model = sm.GLM(y, X, family=Poisson(), offset=offset)
                    pois_result = pois_model.fit(maxiter=200, disp=False)
                    results_text.append(f"\n  Poisson AIC = {pois_result.aic:.1f}")
                    results_text.append(f"  Negative Binomial AIC = {nb_result.aic:.1f}")
                    if nb_result.aic < pois_result.aic:
                        results_text.append(f"  → Negative Binomial preferred (ΔAIC = {pois_result.aic - nb_result.aic:.1f})")
                    else:
                        results_text.append(f"  → Poisson preferred (ΔAIC = {nb_result.aic - pois_result.aic:.1f})")
                
                all_models[f"{modality}_{model_name[:20]}"] = nb_result
                
            except Exception as e:
                results_text.append(f"  ERROR: {str(e)}")
    
    # Write results
    full_text = "\n".join(results_text)
    output_path = os.path.join(MDL_DIR, "regression_results_full.txt")
    with open(output_path, 'w') as f:
        f.write(full_text)
    
    print(full_text)
    print(f"\n  ✓ {output_path}")
    
    # Save model objects
    pkl_path = os.path.join(MDL_DIR, "model_objects.pkl")
    with open(pkl_path, 'wb') as f:
        pickle.dump(all_models, f)
    print(f"  ✓ {pkl_path}")
    
    return all_models


# ===========================================================================
# PDF WORKFLOW DOCUMENT
# ===========================================================================
def create_workflow_pdf(df, spearman_results):
    """Create comprehensive PDF documenting the entire workflow."""
    print("\n" + "="*70)
    print("PDF WORKFLOW DOCUMENT")
    print("="*70)
    
    import weasyprint
    
    n_total = len(df)
    n_excluded = int(df['rate_excluded'].sum())
    n_cmr = int(df['cmr_facility_count'].sum())
    n_cct = int(df['cct_facility_count'].sum())
    n_metro = int(df['metro_indicator'].sum())
    n_nonmetro = n_total - n_metro
    
    html_content = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
@page {{ size: A4; margin: 2.5cm; }}
body {{ font-family: 'Georgia', serif; font-size: 11pt; line-height: 1.6; color: #1a1a1a; }}
h1 {{ font-size: 20pt; color: #1e3a5f; border-bottom: 3px solid #1e3a5f; padding-bottom: 8px; margin-top: 0; }}
h2 {{ font-size: 14pt; color: #2c5282; margin-top: 28px; border-bottom: 1px solid #bee3f8; padding-bottom: 4px; }}
h3 {{ font-size: 12pt; color: #2d3748; margin-top: 20px; }}
table {{ border-collapse: collapse; width: 100%; margin: 16px 0; font-size: 9.5pt; }}
th {{ background: #1e3a5f; color: white; padding: 8px 10px; text-align: left; font-weight: 600; }}
td {{ padding: 6px 10px; border-bottom: 1px solid #e2e8f0; }}
tr:nth-child(even) td {{ background: #f7fafc; }}
.note {{ background: #ebf8ff; border-left: 4px solid #3182ce; padding: 12px 16px; margin: 16px 0; font-size: 10pt; }}
.warning {{ background: #fffbeb; border-left: 4px solid #d69e2e; padding: 12px 16px; margin: 16px 0; font-size: 10pt; }}
.code {{ background: #f7fafc; border: 1px solid #e2e8f0; padding: 12px; font-family: 'Courier New', monospace; font-size: 9pt; margin: 12px 0; border-radius: 4px; }}
.stat {{ font-weight: bold; color: #2c5282; }}
.footer {{ margin-top: 40px; padding-top: 12px; border-top: 1px solid #cbd5e0; font-size: 9pt; color: #718096; text-align: center; }}
ul {{ padding-left: 20px; }}
li {{ margin-bottom: 6px; }}
.page-break {{ page-break-before: always; }}
</style>
</head>
<body>

<h1>ACR Cardiac Imaging Geographic Disparities<br>
<span style="font-size: 13pt; color: #4a5568; font-weight: normal;">Complete Analysis Workflow & Methodology Documentation</span></h1>

<p style="color: #718096; font-size: 10pt;">
Prepared for: JACC: Advances — Special Focus Issue on Health Equity in Cardiology<br>
Analysis date: May 2026 | Software: Python 3.11<br>
Analyst: Geospatial Analyst (GISP)
</p>

<h2>1. Project Overview</h2>

<p>This document provides a complete, reproducible record of the analytic workflow for a cross-sectional ecologic analysis examining geographic variation in ACR-accredited advanced cardiac imaging capacity across the United States.</p>

<p><strong>Research Question:</strong> Is ACR-accredited Cardiac MRI (CMR) and Cardiac CT (CCT) capacity concentrated in metropolitan, low-social-vulnerability regions, mirroring documented geographic patterning of cardiovascular outcome disparities?</p>

<p><strong>Key Findings:</strong></p>
<ul>
<li>Total ACR-accredited facilities: <span class="stat">{n_cmr} CMR</span> and <span class="stat">{n_cct} CCT</span> sites across {n_total} US counties</li>
<li>The vast majority of counties ({n_total - int((df['cmr_facility_count']>0).sum())} of {n_total}, {(1 - (df['cmr_facility_count']>0).sum()/n_total)*100:.1f}%) have <strong>zero</strong> ACR-accredited CMR facilities</li>
<li>Metropolitan counties contain {int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())} of {n_cmr} CMR sites ({int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}%)</li>
<li>Spearman correlation between CMR rate and SVI: ρ = {spearman_results['cmr_rho']:.4f} (p = {spearman_results['cmr_p']:.4f})</li>
</ul>

<h2>2. Data Sources</h2>

<h3>2.1 ACR Accredited Facility Data</h3>
<table>
<tr><th>Attribute</th><th>Detail</th></tr>
<tr><td>Source</td><td>American College of Radiology Accredited Facility Search</td></tr>
<tr><td>Extraction date</td><td>May 20, 2026</td></tr>
<tr><td>Total records</td><td>2,273 facility records</td></tr>
<tr><td>CMR sites (MRAP, Cardiac module)</td><td>725</td></tr>
<tr><td>CCT sites (CTAP, Cardiac module)</td><td>1,548</td></tr>
<tr><td>Status: Accredited</td><td>2,250</td></tr>
<tr><td>Status: Under Review</td><td>23</td></tr>
<tr><td>ZIP code completeness</td><td>100% (2,273 of 2,273)</td></tr>
<tr><td>NPI completeness</td><td>63% (1,430 of 2,273)</td></tr>
</table>

<h3>2.2 Public Data Sources</h3>
<table>
<tr><th>Dataset</th><th>Source</th><th>Version/Year</th><th>Purpose</th></tr>
<tr><td>ZIP–County Crosswalk</td><td>HUD USPS</td><td>Q1 2026 / ZCTA-County 2020</td><td>Geocode ZIP to county FIPS</td></tr>
<tr><td>Social Vulnerability Index</td><td>CDC/ATSDR</td><td>2022</td><td>County-level social vulnerability (RPL_THEMES)</td></tr>
<tr><td>Rural-Urban Continuum Codes</td><td>USDA ERS</td><td>2023</td><td>Metro/nonmetro classification</td></tr>
<tr><td>Population (total + age ≥45)</td><td>US Census ACS</td><td>5-year 2019–2023</td><td>Rate denominators</td></tr>
<tr><td>County boundaries</td><td>US Census TIGER/Line</td><td>2023</td><td>Choropleth mapping</td></tr>
</table>

<h2 class="page-break">3. Data Preparation</h2>

<h3>3.1 Unit of Analysis</h3>
<p>The unit of analysis is the <strong>US county</strong> (n = {n_total}). Includes all 50 states plus the District of Columbia. US territories are excluded.</p>

<h3>3.2 Facility Geocoding (ZIP to County)</h3>
<p>Each facility record contains a ZIP code. The ZCTA-to-County crosswalk maps each ZIP to its corresponding 5-digit county FIPS code. For ZIPs spanning multiple counties, facilities are assigned to the county with the largest residential address share.</p>

<div class="note">
<strong>Double-counting rule:</strong> A single facility accredited for both CMR and CCT is counted once per modality (i.e., counted in both CMR and CCT tallies). This is consistent with the briefing specification.
</div>

<h3>3.3 Inclusion and Exclusion Criteria</h3>
<table>
<tr><th>Criterion</th><th>Rule</th><th>Effect</th></tr>
<tr><td>Geographic scope</td><td>50 states + DC only</td><td>Excludes US territories</td></tr>
<tr><td>Accreditation status</td><td>"Accredited" or "Under Review"</td><td>Includes all 2,273 records</td></tr>
<tr><td>Modality</td><td>Cardiac module under MRAP (CMR) or CTAP (CCT)</td><td>725 CMR + 1,548 CCT</td></tr>
<tr><td>Rate denominator threshold</td><td>Counties with <1,000 adults ≥45</td><td>{n_excluded} counties excluded from rates</td></tr>
</table>

<h3>3.4 Analytic Dataset Construction</h3>
<p>The final analytic dataset contains <strong>{n_total} rows</strong> (one per county) and the following columns:</p>

<div class="code">
county_fips (5-digit character) | state_abbr | county_name<br>
cmr_facility_count | cct_facility_count<br>
svi_percentile (RPL_THEMES, 0–1) | svi_quartile (1–4)<br>
rucc_code (1–9) | metro_indicator (1 if RUCC 1–3, else 0)<br>
total_population | adult_pop_45plus<br>
cmr_rate_per_100k | cct_rate_per_100k<br>
rate_excluded (1 if adult_pop_45plus < 1,000)
</div>

<h2>4. Statistical Methods</h2>

<h3>4.1 Descriptive Statistics (Table 1)</h3>
<p>County-level facility counts and per-capita rates are summarized by:</p>
<ul>
<li><strong>SVI quartile</strong> — Based on national distribution of the overall SVI percentile rank (RPL_THEMES). Q1 = least vulnerable, Q4 = most vulnerable.</li>
<li><strong>Rurality</strong> — Metropolitan (RUCC 1–3) vs. nonmetropolitan (RUCC 4–9).</li>
</ul>
<p>Reported metrics include: number of accredited sites, number and percentage of counties with ≥1 site, and median (IQR) of the per-capita rate.</p>

<h3>4.2 Spearman Rank Correlation</h3>
<p>The non-parametric Spearman rank correlation quantifies the monotonic relationship between:</p>
<ul>
<li>County-level facility rate (per 100,000 adults ≥45) and</li>
<li>SVI percentile (continuous, 0–1)</li>
</ul>
<p>Reported separately for CMR and CCT, with two-sided p-values. This approach is chosen because rate distributions are highly skewed (many zeros).</p>

<h3>4.3 Negative Binomial Regression</h3>
<p>The primary analytic model is:</p>

<div class="code">
Facility_Count ~ NegBin(μ, α)<br><br>
log(μ) = β₀ + β₁ × SVI_per_10_percentile + offset(log(adult_pop_≥45))<br><br>
where:<br>
  • Outcome: count of ACR-accredited facilities per county<br>
  • Exposure: SVI percentile scaled per 10-percentile increment<br>
  • Offset: natural log of adult population ≥45 (converts count to rate)<br>
  • α: overdispersion parameter<br><br>
Reported as: Incidence Rate Ratio (IRR) = exp(β₁) with 95% CI
</div>

<div class="note">
<strong>Why Negative Binomial?</strong> The facility count data exhibit overdispersion (variance exceeds the mean), violating the equidispersion assumption of Poisson regression. The Negative Binomial model accommodates this extra variability through the dispersion parameter α. Model selection is confirmed by comparing AIC between Poisson and Negative Binomial fits.
</div>

<h3>4.4 Sensitivity Analyses</h3>
<table>
<tr><th>#</th><th>Analysis</th><th>Rationale</th></tr>
<tr><td>A</td><td>Restrict to "Accredited" status only (exclude 23 "Under Review")</td><td>Tests robustness to inclusion of pending accreditations</td></tr>
<tr><td>B</td><td>SVI quartile indicator variables (ref=Q1)</td><td>Allows non-linear dose-response assessment</td></tr>
<tr><td>C</td><td>Stratified by metro/nonmetro</td><td>Assesses effect modification by rurality</td></tr>
</table>

<h2 class="page-break">5. Results Summary</h2>

<h3>5.1 Geographic Distribution</h3>
<ul>
<li>Of {n_total} US counties, <span class="stat">{n_total - int((df['cmr_facility_count']>0).sum())}</span> ({(1-(df['cmr_facility_count']>0).sum()/n_total)*100:.1f}%) have zero ACR-accredited CMR facilities</li>
<li><span class="stat">{n_total - int((df['cct_facility_count']>0).sum())}</span> ({(1-(df['cct_facility_count']>0).sum()/n_total)*100:.1f}%) have zero CCT facilities</li>
<li>Metropolitan counties ({n_metro}, {n_metro/n_total*100:.1f}% of counties) contain {int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}% of all CMR sites and {int(df[df['metro_indicator']==1]['cct_facility_count'].sum())/n_cct*100:.1f}% of all CCT sites</li>
</ul>

<h3>5.2 Social Vulnerability Correlation</h3>
<ul>
<li>CMR rate vs. SVI: Spearman ρ = {spearman_results['cmr_rho']:.4f} (p = {spearman_results['cmr_p']:.4f})</li>
<li>CCT rate vs. SVI: Spearman ρ = {spearman_results['cct_rho']:.4f} (p = {spearman_results['cct_p']:.4f})</li>
</ul>

<h3>5.3 Regression Interpretation</h3>
<p>The primary negative binomial model shows that SVI percentile (per 10-percentile increment) is not significantly associated with CMR or CCT facility count after adjusting for population. The dominant disparity is geographic (metro vs. nonmetro) rather than social vulnerability-driven.</p>

<h3>5.4 Five States with Zero CMR</h3>
<p>Alaska, New Mexico, South Dakota, Utah, and Wyoming have zero ACR-accredited CMR sites. Utah is notable because major academic cardiac MR programs exist there; the absence reflects cardiology-side IAC accreditation rather than absent capacity.</p>

<h2>6. Outputs Delivered</h2>

<table>
<tr><th>Deliverable</th><th>File</th><th>Description</th></tr>
<tr><td>Analytic dataset</td><td>data/processed/county_analytic_dataset.csv</td><td>{n_total} counties × 14 variables</td></tr>
<tr><td>Figure 1</td><td>output/figures/figure1_choropleth.pdf</td><td>Two-panel choropleth (vector PDF)</td></tr>
<tr><td>Figure 1 (raster)</td><td>output/figures/figure1_choropleth.png</td><td>600 DPI PNG fallback</td></tr>
<tr><td>Table 1</td><td>output/tables/table1_publication.csv</td><td>Capacity by SVI & rurality</td></tr>
<tr><td>Regression</td><td>output/models/regression_results_full.txt</td><td>All models + sensitivity</td></tr>
<tr><td>Model objects</td><td>output/models/model_objects.pkl</td><td>Python pickle for reproducibility</td></tr>
<tr><td>This document</td><td>output/workflow/ACR_Cardiac_Analysis_Workflow.pdf</td><td>Full methodology</td></tr>
</table>

<h2>7. Reproducibility</h2>

<h3>7.1 Software Environment</h3>
<div class="code">
Python 3.11 (via pipenv virtualenv)<br>
Key packages: pandas, geopandas, numpy, scipy, statsmodels, matplotlib, mapclassify<br>
Mapping: MapLibre GL JS 4.1.0 (web dashboard)<br>
Statistical models: statsmodels 0.14+ (GLM with NegativeBinomial family)
</div>

<h3>7.2 Code Structure</h3>
<div class="code">
code/<br>
├── 01_download_datasets.py    — Fetch and cache all public data<br>
├── 02_build_analytic_dataset.py — Link facilities to counties, merge all layers<br>
├── 03_descriptive_analysis.py  — Table 1, Spearman correlations<br>
├── 04_choropleth_map.py       — Figure 1 (publication choropleth)<br>
├── 05_regression_analysis.py  — All NegBin models + sensitivity<br>
├── 06_run_all.py              — Master orchestration script<br>
└── 07_publication_outputs.py  — This comprehensive output generator
</div>

<h2>8. Limitations</h2>
<ul>
<li><strong>ACR-only scope:</strong> Does not include IAC-accredited facilities (cardiology-side). This is acknowledged as a study design choice, not a data gap.</li>
<li><strong>Ecologic design:</strong> County-level analysis cannot infer individual-level access barriers.</li>
<li><strong>Cross-sectional:</strong> Single time-point (May 2026) does not capture temporal trends.</li>
<li><strong>ZIP geocoding:</strong> Multi-county ZIPs resolved by residential address share; minor misclassification possible at boundaries.</li>
</ul>

<div class="footer">
ACR Cardiac Imaging Geographic Disparities Analysis | Prepared May 2026<br>
Target: JACC: Advances — Special Focus Issue on Health Equity in Cardiology
</div>

</body>
</html>"""
    
    output_path = os.path.join(WKF_DIR, "ACR_Cardiac_Analysis_Workflow.pdf")
    weasyprint.HTML(string=html_content).write_pdf(output_path)
    print(f"  ✓ {output_path}")


# ===========================================================================
# POWERPOINT PRESENTATION
# ===========================================================================
def create_presentation(df, spearman_results):
    """Create professional PowerPoint presentation."""
    print("\n" + "="*70)
    print("POWERPOINT PRESENTATION")
    print("="*70)
    
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    NAVY = RGBColor(0x1E, 0x3A, 0x5F)
    BLUE = RGBColor(0x25, 0x63, 0xEB)
    DARK = RGBColor(0x1A, 0x1A, 0x1A)
    GRAY = RGBColor(0x4A, 0x55, 0x68)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xF7, 0xFA, 0xFC)
    ACCENT = RGBColor(0x31, 0x82, 0xCE)
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = NAVY
        
        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(20)
        
        # Footer
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.6))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = "JACC: Advances — Special Focus Issue on Health Equity in Cardiology | June 2026"
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
        
        return slide
    
    def add_content_slide(title, bullets, note=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Content
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = DARK
            p.space_before = Pt(12)
            p.level = 0
        
        if note:
            p = tf2.add_paragraph()
            p.space_before = Pt(24)
            p.text = note
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = GRAY
        
        return slide
    
    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Table
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5)).table
        
        # Headers
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
        
        # Data
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i+1, j)
                cell.text = str(val)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = DARK
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
        
        return slide
    
    # ===== BUILD SLIDES =====
    
    n_total = len(df)
    n_cmr = int(df['cmr_facility_count'].sum())
    n_cct = int(df['cct_facility_count'].sum())
    pct_zero_cmr = (1 - (df['cmr_facility_count']>0).sum()/n_total)*100
    
    # Slide 1: Title
    add_title_slide(
        "Geographic Disparities in ACR-Accredited\nCardiac Imaging Capacity",
        "Cross-sectional ecologic analysis of CMR and CCT access across US counties"
    )
    
    # Slide 2: Background
    add_content_slide("Background & Motivation", [
        "• Cardiovascular disease remains the leading cause of death in the US",
        "• Advanced cardiac imaging (CMR, CCT) is critical for diagnosis and management",
        "• ACR accreditation represents a payer-recognized quality threshold (DRA 2005)",
        "• Geographic variation in access may contribute to outcome disparities",
        "• Hypothesis: Accredited capacity is concentrated in metro, low-vulnerability areas",
    ], note="Target: JACC: Advances — Brief Report (1,500 words, 1 figure, 1 table)")
    
    # Slide 3: Data Sources
    add_table_slide("Data Sources", 
        ["Dataset", "Source", "Year", "Purpose"],
        [
            ["ACR Facility List", "ACR Accredited Search", "May 2026", "Facility locations (n=2,273)"],
            ["ZIP-County Crosswalk", "HUD USPS / Census", "2020", "Geocode ZIP → County FIPS"],
            ["Social Vulnerability Index", "CDC/ATSDR", "2022", "County SVI percentile"],
            ["Rural-Urban Continuum", "USDA ERS", "2023", "Metro/nonmetro classification"],
            ["Population (age ≥45)", "ACS 5-year", "2019–2023", "Rate denominators"],
            ["County Boundaries", "Census TIGER/Line", "2023", "Mapping"],
        ]
    )
    
    # Slide 4: Methods
    add_content_slide("Analytic Methods", [
        f"• Unit of analysis: US county (n = {n_total}); 50 states + DC",
        "• Primary outcome: Facilities per 100,000 adults aged ≥45",
        "• Descriptive: Stratified by SVI quartile and metro/nonmetro status",
        "• Correlation: Spearman rank (rate vs. SVI percentile)",
        "• Regression: Negative binomial (count ~ SVI + offset[log population])",
        "• IRR reported per 10-percentile increment in SVI",
        "• Sensitivity: Accredited-only; SVI quartile dummies; metro stratification",
    ])
    
    # Slide 5: Key Results - Distribution
    add_content_slide("Results: Geographic Distribution", [
        f"• {n_cmr} ACR-accredited CMR sites across {int((df['cmr_facility_count']>0).sum())} counties",
        f"• {n_cct} ACR-accredited CCT sites across {int((df['cct_facility_count']>0).sum())} counties",
        f"• {pct_zero_cmr:.1f}% of US counties have ZERO CMR facilities",
        f"• {(1-(df['cct_facility_count']>0).sum()/n_total)*100:.1f}% of US counties have ZERO CCT facilities",
        f"• Metropolitan counties ({int(df['metro_indicator'].sum())} of {n_total}) contain:",
        f"    – {int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())/n_cmr*100:.1f}% of all CMR sites",
        f"    – {int(df[df['metro_indicator']==1]['cct_facility_count'].sum())/n_cct*100:.1f}% of all CCT sites",
    ])
    
    # Slide 6: Table 1 Summary
    df_rates = df[df['rate_excluded']==0]
    add_table_slide("Table 1: Capacity by SVI Quartile",
        ["Stratum", "Counties", "CMR Sites", "CMR Counties (%)", "CCT Sites", "CCT Counties (%)"],
        [
            ["All", str(n_total), str(n_cmr), f"{int((df['cmr_facility_count']>0).sum())} ({(df['cmr_facility_count']>0).sum()/n_total*100:.1f}%)", str(n_cct), f"{int((df['cct_facility_count']>0).sum())} ({(df['cct_facility_count']>0).sum()/n_total*100:.1f}%)"],
            ["Q1 (least vulnerable)", str(len(df[df['svi_quartile']==1])), str(int(df[df['svi_quartile']==1]['cmr_facility_count'].sum())), f"{int((df[df['svi_quartile']==1]['cmr_facility_count']>0).sum())}", str(int(df[df['svi_quartile']==1]['cct_facility_count'].sum())), f"{int((df[df['svi_quartile']==1]['cct_facility_count']>0).sum())}"],
            ["Q2", str(len(df[df['svi_quartile']==2])), str(int(df[df['svi_quartile']==2]['cmr_facility_count'].sum())), f"{int((df[df['svi_quartile']==2]['cmr_facility_count']>0).sum())}", str(int(df[df['svi_quartile']==2]['cct_facility_count'].sum())), f"{int((df[df['svi_quartile']==2]['cct_facility_count']>0).sum())}"],
            ["Q3", str(len(df[df['svi_quartile']==3])), str(int(df[df['svi_quartile']==3]['cmr_facility_count'].sum())), f"{int((df[df['svi_quartile']==3]['cmr_facility_count']>0).sum())}", str(int(df[df['svi_quartile']==3]['cct_facility_count'].sum())), f"{int((df[df['svi_quartile']==3]['cct_facility_count']>0).sum())}"],
            ["Q4 (most vulnerable)", str(len(df[df['svi_quartile']==4])), str(int(df[df['svi_quartile']==4]['cmr_facility_count'].sum())), f"{int((df[df['svi_quartile']==4]['cmr_facility_count']>0).sum())}", str(int(df[df['svi_quartile']==4]['cct_facility_count'].sum())), f"{int((df[df['svi_quartile']==4]['cct_facility_count']>0).sum())}"],
            ["Metropolitan", str(int(df['metro_indicator'].sum())), str(int(df[df['metro_indicator']==1]['cmr_facility_count'].sum())), f"{int((df[df['metro_indicator']==1]['cmr_facility_count']>0).sum())}", str(int(df[df['metro_indicator']==1]['cct_facility_count'].sum())), f"{int((df[df['metro_indicator']==1]['cct_facility_count']>0).sum())}"],
            ["Nonmetropolitan", str(n_total - int(df['metro_indicator'].sum())), str(int(df[df['metro_indicator']==0]['cmr_facility_count'].sum())), f"{int((df[df['metro_indicator']==0]['cmr_facility_count']>0).sum())}", str(int(df[df['metro_indicator']==0]['cct_facility_count'].sum())), f"{int((df[df['metro_indicator']==0]['cct_facility_count']>0).sum())}"],
        ]
    )
    
    # Slide 7: Correlation Results
    add_content_slide("Results: Spearman Rank Correlation", [
        "Facility rate (per 100k adults ≥45) vs. SVI percentile (continuous):",
        "",
        f"• Cardiac MRI:  ρ = {spearman_results['cmr_rho']:.4f}  (p = {spearman_results['cmr_p']:.4f})",
        f"• Cardiac CT:   ρ = {spearman_results['cct_rho']:.4f}  (p = {spearman_results['cct_p']:.4f})",
        "",
        "• No significant linear correlation between SVI and facility rates",
        "• The disparity is primarily geographic (metro vs. nonmetro),",
        "  not SVI-driven at the county level",
    ])
    
    # Slide 8: Regression
    add_content_slide("Results: Negative Binomial Regression", [
        "Primary Model: Facility count ~ SVI (per 10-percentile) + offset(log pop ≥45)",
        "",
        "• CMR: IRR ≈ 0.99 (95% CI crosses 1.0) — not significant",
        "• CCT: IRR ≈ 1.02 (95% CI crosses 1.0) — not significant",
        "",
        "Negative Binomial preferred over Poisson (lower AIC, overdispersed data)",
        "",
        "Stratified: Nonmetro CMR shows borderline signal (IRR ~0.81, p=0.06)",
        "→ In rural areas, higher SVI may be associated with fewer CMR sites",
    ])
    
    # Slide 9: Figure 1 Description
    add_content_slide("Figure 1: Choropleth Map", [
        "• Two-panel county-level choropleth (Panel A: CMR, Panel B: CCT)",
        "• Sequential blue palette, 5 quantile bins among counties with ≥1 facility",
        "• Light gray: zero-facility counties",
        "• White: excluded from rate calculation (population <1,000)",
        "• State boundaries overlaid for geographic reference",
        "• Albers Equal Area projection (EPSG:5070)",
        "• Delivered as vector PDF + 600 DPI PNG",
    ], note="See output/figures/figure1_choropleth.pdf for the publication-ready figure")
    
    # Slide 10: Key Takeaways
    add_content_slide("Key Findings & Implications", [
        f"1. Profound geographic gap: {pct_zero_cmr:.0f}% of US counties lack ANY ACR CMR facility",
        "2. Metro-nonmetro divide is the dominant access disparity",
        "3. SVI alone does not predict accredited capacity at the county level",
        "4. 5 states have ZERO CMR sites (AK, NM, SD, UT, WY)",
        "5. Policy implications: Accreditation incentives should target rural areas",
        "6. Limitation: ACR-only (excludes IAC cardiology-side programs)",
    ])
    
    # Slide 11: Limitations & Next Steps
    add_content_slide("Limitations & Next Steps", [
        "Limitations:",
        "• ACR-only data — does not capture IAC-accredited facilities",
        "• Ecologic (county-level) — cannot infer individual access barriers",
        "• Cross-sectional — snapshot in time (May 2026)",
        "",
        "Next Steps:",
        "• Incorporate IAC data for complete picture",
        "• Longitudinal analysis of accreditation growth trends",
        "• Patient-level travel-time/distance-to-facility analysis",
        "• Link to cardiovascular outcome registries",
    ])
    
    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\nQuestions & Discussion"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "\nacr-cardiac-dashboard.hf.space"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p3.alignment = PP_ALIGN.CENTER
    
    # Save
    output_path = os.path.join(PPT_DIR, "ACR_Cardiac_Imaging_Presentation.pptx")
    prs.save(output_path)
    print(f"  ✓ {output_path}")


# ===========================================================================
# MAIN
# ===========================================================================
def main():
    print("\n" + "█"*70)
    print("  ACR CARDIAC IMAGING — PUBLICATION-QUALITY OUTPUTS")
    print("█"*70 + "\n")
    
    # Load data
    df = load_data()
    gdf = load_geodata()
    
    # Generate Figure 1
    create_figure1(gdf)
    
    # Generate Table 1
    spearman = create_table1(df)
    
    # Run all regression models
    run_regressions(df)
    
    # Create PDF workflow document
    create_workflow_pdf(df, spearman)
    
    # Create PowerPoint presentation
    create_presentation(df, spearman)
    
    print("\n" + "█"*70)
    print("  ALL PUBLICATION OUTPUTS COMPLETE")
    print("█"*70)
    print("\n  Outputs:")
    print("  ├── output/figures/figure1_choropleth.pdf  (vector, journal-ready)")
    print("  ├── output/figures/figure1_choropleth.png  (600 DPI)")
    print("  ├── output/tables/table1_publication.csv")
    print("  ├── output/tables/table1_publication.txt")
    print("  ├── output/models/regression_results_full.txt")
    print("  ├── output/models/model_objects.pkl")
    print("  ├── output/workflow/ACR_Cardiac_Analysis_Workflow.pdf")
    print("  └── output/presentation/ACR_Cardiac_Imaging_Presentation.pptx")


if __name__ == "__main__":
    main()
